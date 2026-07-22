"""TimeHeroes — Cours App Mobile V2 (estimations Ronald + explications débutant)"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── PALETTE ──────────────────────────────────────────────────────────
BG = RGBColor(0xF8, 0xF5, 0xF0)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2 = RGBColor(0xF0, 0xED, 0xE8)
ACCENT = RGBColor(0x2B, 0xB2, 0x86)
ACCENT_DARK = RGBColor(0x00, 0x8F, 0x78)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_SEC = RGBColor(0x6B, 0x6B, 0x6B)
TEXT_MUTED = RGBColor(0x9A, 0x9A, 0x9A)
BORDER = RGBColor(0xE5, 0xE0, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN_DARK = RGBColor(0x06, 0x78, 0x54)

W = Inches(13.33)
H = Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color=WHITE, border=None, radius=None):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(st, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if radius:
        s.adjustments[0] = radius
    return s

def circle(slide, l, t, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def tb(slide, l, t, w, h):
    tbox = slide.shapes.add_textbox(l, t, w, h)
    tbox.text_frame.word_wrap = True
    return tbox.text_frame

def P(tf, text="", size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, name="Calibri"):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    p.space_after = Pt(4)
    return p

def new_slide():
    s = prs.slides.add_slide(BLANK)
    set_bg(s, BG)
    return s

def card(slide, l, t, w, h):
    return rect(slide, l, t, w, h, SURFACE, BORDER, 0.02)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 - TITLE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, ACCENT_DARK)
circle(s, Inches(-1), Inches(-1), Inches(3), RGBColor(0x22, 0xCC, 0xA0))
circle(s, Inches(11), Inches(5), Inches(4), RGBColor(0x00, 0xAA, 0x88))
circle(s, Inches(10), Inches(-2), Inches(2.5), RGBColor(0x00, 0x99, 0x77))

tf = tb(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1))
P(tf, "📱 Transformer TimeHeroes en App Mobile", 40, True, WHITE, PP_ALIGN.LEFT, "Calibri Light")
tf = tb(s, Inches(1.5), Inches(3), Inches(10), Inches(1.5))
P(tf, "Cours complet pour les non-initiés", 18, False, RGBColor(0xCC, 0xFA, 0xE8))
P(tf, "PWA · Capacitor · React Native · les 3 options expliquées pas à pas", 14, False, RGBColor(0xAA, 0xEE, 0xDD))
tf = tb(s, Inches(1.5), Inches(5.5), Inches(5), Inches(1))
P(tf, "Projet : TimeHeroes — Banque du Temps", 12, False, RGBColor(0xAA, 0xEE, 0xDD))
P(tf, "Stack : Next.js · Prisma · SQLite", 12, False, RGBColor(0xAA, 0xEE, 0xDD))
P(tf, "Fait par Ronald (avec Hermes Agent) — Juillet 2026", 12, False, RGBColor(0xAA, 0xEE, 0xDD))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 - C'EST QUOI UNE APP MOBILE ? (concepts de base)
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "🤔 C'est quoi une app mobile ?", 30, True, TEXT, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5))
P(tf, "Les bases pour comprendre la suite — sans jargon", 14, False, TEXT_SEC)

# 3 concepts
concepts = [
    ("🌐", "App Web (actuelle)", "C'est TimeHeroes aujourd'hui. Tu ouvres Chrome ou Safari sur ton téléphone, tu tapes timeheroes.fr, et le site s'affiche. Pour y accéder, tu as BESOIN d'Internet et d'un navigateur.", ACCENT),
    ("📲", "App installable (PWA)", "C'est le MÊME site web, mais le navigateur te propose de l'\"installer\" sur ton écran d'accueil. Elle a sa propre icône et peut fonctionner un peu hors-ligne. Pas besoin d'App Store.", PURPLE),
    ("📱", "App native", "C'est comme Instagram, Uber, WhatsApp. Tu la télécharges depuis l'App Store (iPhone) ou Google Play (Android). Elle peut tout faire : caméra, NFC, notifications, Bluetooth. C'est le plus puissant mais le plus long à fabriquer.", AMBER),
]
for i, (emoji, title, desc, color) in enumerate(concepts):
    y = Inches(1.9) + Inches(i * 1.7)
    card(s, Inches(0.8), y, Inches(11.5), Inches(1.5))
    circle(s, Inches(1.1), y + Inches(0.2), Inches(1.1), color)
    tf = tb(s, Inches(1.1), y + Inches(0.3), Inches(1.1), Inches(0.8))
    P(tf, emoji, 28, False, WHITE, PP_ALIGN.CENTER)
    tf = tb(s, Inches(2.5), y + Inches(0.2), Inches(3), Inches(0.5))
    P(tf, title, 18, True, color, name="Calibri Light")
    tf = tb(s, Inches(2.5), y + Inches(0.65), Inches(9.3), Inches(0.7))
    P(tf, desc, 12, False, TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 - POURQUOI PASSER EN MOBILE ?
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "🎯 Pourquoi ne pas rester sur le site web ?", 30, True, TEXT, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5))
P(tf, "Les super-pouvoirs que tu gagnes en passant en mobile", 14, False, TEXT_SEC)

data = [
    ("📸 Scanner un QR code", "Le site peut utiliser la caméra, mais c'est limité. Une app native scanne bien plus vite et mieux.", "✅ Scan natif"),
    ("🤝 Valider par NFC (sans contact)", "Impossible depuis un site web. Le NFC nécessite une app native. Avec une app, tu valides une mission en tapant les téléphones.", "✅ Tap NFC !"),
    ("🔔 Recevoir des notifications", "Sur un site, tu reçois une notification uniquement si tu as l'onglet ouvert. En app, tu reçois tout le temps.", "✅ Push toujours"),
    ("📴 Utiliser sans Internet", "Le site ne marche pas sans connexion. Une app peut mettre en cache les données pour les consulter hors-ligne.", "✅ Mode hors-ligne"),
    ("🏠 Avoir une icône sur le téléphone", "Un site peut être ajouté à l'écran d'accueil, mais c'est moins bien qu'une vraie app installée depuis un store.", "✅ Icône permanente"),
]
for i, (label, expl, win) in enumerate(data):
    y = Inches(1.8) + Inches(i * 1.05)
    card(s, Inches(0.8), y, Inches(7.5), Inches(0.9))
    tf = tb(s, Inches(1), y + Inches(0.08), Inches(7), Inches(0.3))
    P(tf, label, 14, True, TEXT)
    tf = tb(s, Inches(1), y + Inches(0.4), Inches(7), Inches(0.4))
    P(tf, expl, 10, False, TEXT_SEC)
    rect(s, Inches(8.5), y, Inches(3.8), Inches(0.9), RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0x6E, 0xE7, 0xB7), 0.02)
    tf = tb(s, Inches(8.6), y + Inches(0.2), Inches(3.6), Inches(0.5))
    P(tf, win, 12, True, ACCENT_DARK, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 - LES 3 OPTIONS (VUE D'ENSEMBLE)
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "🛤️ Les 3 options pour rendre TimeHeroes mobile", 30, True, TEXT, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4))
P(tf, "Du plus simple au plus puissant — chaque option est expliquée en détail après", 14, False, TEXT_SEC)

opts = [
    ("Option 1 : PWA", "1 jour", ACCENT, "✅ Gratuit, facile", "❌ Pas de NFC"),
    ("Option 2 : Capacitor", "3-4 jours", PURPLE, "✅ NFC ! ✅ 95% du code gardé", "❌ Compte Apple 99€/an"),
    ("Option 3 : React Native", "2-4 semaines", AMBER, "✅ Tout natif, pro", "❌ Cher, long, code à réécrire"),
]
for i, (name, time, color, pro, con) in enumerate(opts):
    x = Inches(0.8) + Inches(i * 4.2)
    card(s, x, Inches(1.8), Inches(3.8), Inches(5.2))
    rect(s, x, Inches(1.8), Inches(3.8), Inches(0.06), color)
    circle(s, x + Inches(1.4), Inches(2.1), Inches(1), color)
    tf = tb(s, x + Inches(1.4), Inches(2.1), Inches(1), Inches(1))
    P(tf, str(i+1), 32, True, WHITE, PP_ALIGN.CENTER, "Calibri Light")
    tf = tb(s, x + Inches(0.2), Inches(3.3), Inches(3.4), Inches(0.5))
    P(tf, name, 20, True, TEXT, PP_ALIGN.CENTER, "Calibri Light")
    tf = tb(s, x + Inches(0.2), Inches(3.8), Inches(3.4), Inches(0.4))
    P(tf, time, 13, True, color, PP_ALIGN.CENTER)
    tf = tb(s, x + Inches(0.3), Inches(4.4), Inches(3.2), Inches(0.5))
    P(tf, pro, 11, False, GREEN_DARK)
    tf = tb(s, x + Inches(0.3), Inches(5.0), Inches(3.2), Inches(0.5))
    P(tf, con, 11, False, RGBColor(0xDC, 0x26, 0x26))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 - PWA (EXPLIQUÉ SIMPLEMENT)
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, ACCENT)
tf = tb(s, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7))
P(tf, "🌐 Option 1 : PWA (Progressive Web App)", 32, True, WHITE, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(12), Inches(0.8))
P(tf, "C'est quoi ?", 16, True, RGBColor(0xCC, 0xFA, 0xE8))
P(tf, "Imagine que TimeHeroes reste un site web normal, mais que le navigateur (Chrome, Safari) propose de l'\"installer\" sur ton téléphone.", 12, False, RGBColor(0xE0, 0xF7, 0xEE))
P(tf, "L'utilisateur ouvre timeheroes.fr → Chrome lui dit \"Ajouter à l'écran d'accueil\" → une icône apparaît → au clic, ça s'ouvre comme une vraie app (sans la barre d'adresse). Ça peut même marcher un peu hors-ligne.", 12, False, RGBColor(0xE0, 0xF7, 0xEE))

# Comment ça marche
tf = tb(s, Inches(0.8), Inches(3.2), Inches(12), Inches(0.4))
P(tf, "Comment ça marche techniquement ? (simplifié)", 16, True, RGBColor(0xCC, 0xFA, 0xE8))

steps = [
    ("📄 manifest.json", "Un fichier tout simple qui dit \"voici le nom de l'app, son icône, sa couleur\". Le navigateur lit ce fichier pour savoir comment afficher l'app installée."),
    ("⚙️ Service Worker", "Un petit programme JavaScript qui tourne en arrière-plan (même quand la page est fermée). Il s'occupe de : 1) Cacher des pages pour le hors-ligne 2) Recevoir les notifications push"),
    ("🔔 Web Push API", "C'est le système de notifications. Le serveur envoie un message → ton téléphone affiche une notification. Pas besoin d'App Store."),
    ("📱 Installation", "Sur Android : Chrome propose tout seul l'installation. Sur iPhone : l'utilisateur doit aller dans \"Partager\" → \"Sur l'écran d'accueil\""),
]
for i, (title, desc) in enumerate(steps):
    y = Inches(3.8) + Inches(i * 0.85)
    card(s, Inches(0.8), y, Inches(11.5), Inches(0.75))
    tf = tb(s, Inches(1.1), y + Inches(0.05), Inches(11), Inches(0.65))
    P(tf, title, 13, True, WHITE)
    P(tf, desc, 10, False, TEXT_SEC)

# Bottom bar
rect(s, Inches(0.8), Inches(7), Inches(11.5), Inches(0.5), RGBColor(0xEC, 0xFD, 0xF5), ACCENT, 0.02)
tf = tb(s, Inches(1), Inches(7), Inches(11), Inches(0.5))
P(tf, "⏱️ Estimation si Ronald fait tout : 1 journée  •  💰 Coût : 0€  •  📱 Exemples : Twitter, Pinterest, Starbucks, Uber", 12, True, ACCENT_DARK, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 - CE QUE LA PWA PEUT / NE PEUT PAS FAIRE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "✅ PWA : ce qui marche et ce qui ne marche pas", 30, True, TEXT, name="Calibri Light")

ok = [
    ("✅ Icône sur l'écran d'accueil", "Oui, comme une vraie app"),
    ("✅ Marche sans barre de navigateur", "Oui, plein écran"),
    ("✅ Fonctionne hors-ligne (pages en cache)", "Oui, les pages déjà visitées restent accessibles"),
    ("✅ Notifications push", "Oui, même si l'app est fermée"),
    ("✅ Partage (Web Share API)", "Oui, peut partager du contenu via WhatsApp etc."),
    ("⚠️ QR Code via la caméra", "Possible, mais moins fluide que du natif"),
]
nok = [
    ("❌ NFC (sans contact)", "IMPOSSIBLE. Le NFC est un composant physique du téléphone, le site web n'y a pas accès. C'est le point bloquant pour TimeHeroes."),
    ("❌ Bluetooth avancé", "Pas possible non plus"),
    ("❌ Stockage limité", "iPhone limite le cache à ~50 MB"),
    ("❌ Installation sur iPhone", "Pas automatique : l'utilisateur doit savoir le faire"),
    ("❌ Pas dans l'App Store", "Moins crédible pour une fondation / un partenaire"),
]

for i, (label, desc) in enumerate(ok):
    y = Inches(1.3) + Inches(i * 0.55)
    rect(s, Inches(0.8), y, Inches(6), Inches(0.45), RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0x6E, 0xE7, 0xB7), 0.02)
    tf = tb(s, Inches(1), y + Inches(0.03), Inches(5.6), Inches(0.38))
    P(tf, f"{label} — {desc}", 10, False, ACCENT_DARK)

for i, (label, desc) in enumerate(nok):
    y = Inches(1.3) + Inches(i * 0.9)
    rect(s, Inches(7.2), y, Inches(5.5), Inches(0.8), RGBColor(0xFE, 0xF0, 0xF0), RGBColor(0xFE, 0xCB, 0xCB), 0.02)
    tf = tb(s, Inches(7.4), y + Inches(0.05), Inches(5.1), Inches(0.7))
    P(tf, label, 10, True, RGBColor(0xB9, 0x1C, 0x1C))
    P(tf, desc, 9, False, RGBColor(0x8B, 0x1C, 0x1C))

# Warning box
rect(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.8), RGBColor(0xFF, 0xF3, 0xCD), AMBER, 0.02)
tf = tb(s, Inches(1.1), Inches(6.55), Inches(11.2), Inches(0.7))
P(tf, "⚠️ Pour TimeHeroes : la PWA est un bon début (démo, 1 jour de travail), mais si tu veux le NFC pour valider les missions, il faut passer à l'Option 2 (Capacitor).", 11, True, RGBColor(0x85, 0x64, 0x04))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 - CAPACITOR (EXPLIQUÉ SIMPLEMENT)
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, PURPLE)
tf = tb(s, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7))
P(tf, "🔌 Option 2 : Capacitor", 32, True, WHITE, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(12), Inches(1))
P(tf, "C'est quoi ?", 16, True, RGBColor(0xDD, 0xD6, 0xFC))
P(tf, "Capacitor, c'est une \"coquille\" d'app native qui contient ton site web à l'intérieur.", 12, False, RGBColor(0xE8, 0xE0, 0xFC))
P(tf, "Concrètement : tu prends TimeHeroes (Next.js) tel quel, tu l'emballes dans Capacitor, et tu obtiens une vraie app iPhone/Android. La différence avec la PWA : cette coquille peut communiquer avec les composants physiques du téléphone (NFC, caméra, Bluetooth, notifications push natives).", 12, False, RGBColor(0xE8, 0xE0, 0xFC))

tf = tb(s, Inches(0.8), Inches(3.2), Inches(4), Inches(0.4))
P(tf, "Comment ça marche ? (visuellement)", 16, True, RGBColor(0xDD, 0xD6, 0xFC))

# Architecture visuelle
rect(s, Inches(0.8), Inches(3.8), Inches(11.5), Inches(2.5), RGBColor(0xE8, 0xE0, 0xFC), RGBColor(0xBB, 0xA0, 0xF0), 0.02)

# Layer 1 - WebView
rect(s, Inches(1.2), Inches(3.9), Inches(10.7), Inches(0.7), RGBColor(0xE0, 0xD8, 0xF8))
tf = tb(s, Inches(1.5), Inches(3.95), Inches(10), Inches(0.6))
P(tf, "🌐 COUCHE WEB (Next.js inchangé) — timeheroes.fr/dashboard, /services, /bookings…", 12, True, WHITE)
P(tf, "Ton code actuel ne change PAS. Tout ce que tu as déjà construit reste identique.", 10, False, RGBColor(0xCC, 0xC0, 0xF5))

# Layer 2 - Capacitor
rect(s, Inches(1.2), Inches(4.7), Inches(10.7), Inches(0.7), RGBColor(0xD8, 0xD0, 0xF0))
tf = tb(s, Inches(1.5), Inches(4.75), Inches(10), Inches(0.6))
P(tf, "🔌 PONT CAPACITOR — fait le lien entre le site web et le téléphone", 12, True, RGBColor(0xDD, 0xD6, 0xFC))
P(tf, "Quand tu veux utiliser le NFC : ton site appelle un \"plugin\" Capacitor → Capacitor parle au chip NFC du téléphone.", 10, False, RGBColor(0xCC, 0xC0, 0xF5))

# Layer 3 - Native
rect(s, Inches(1.2), Inches(5.5), Inches(10.7), Inches(0.6), RGBColor(0xD0, 0xC8, 0xE8))
tf = tb(s, Inches(1.5), Inches(5.55), Inches(10), Inches(0.5))
P(tf, "📱 COUCHE NATIVE — ton app publiée sur l'App Store et Google Play", 12, True, WHITE)
P(tf, "Les utilisateurs la téléchargent comme Instagram ou WhatsApp.", 10, False, RGBColor(0xCC, 0xC0, 0xF5))

rect(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.6), RGBColor(0xEC, 0xFD, 0xF5), ACCENT, 0.02)
tf = tb(s, Inches(1), Inches(6.7), Inches(11), Inches(0.6))
P(tf, "⏱️ Estimation si Ronald fait tout : 3-4 jours  •  💰 Coûts : 99€ compte Apple (une fois) + 0,30€/tag NFC  •  📱 NFC !", 12, True, PURPLE, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 - CAPACITOR EN DÉTAIL : CE QUE TU DOIS FAIRE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "🔧 Option 2 (suite) : Les étapes concrètes", 30, True, TEXT, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4))
P(tf, "Ce que Ronald doit faire, étape par étape", 14, False, TEXT_SEC)

steps = [
    ("1️⃣ Installer Capacitor", "npm install @capacitor/core @capacitor/cli", "5 min", "Une commande dans le terminal. Capacitor s'installe comme n'importe quelle librairie JavaScript."),
    ("2️⃣ Initialiser le projet", "npx cap init TimeHeroes com.timeheroes.app", "2 min", "Capacitor crée les dossiers de base pour iOS et Android."),
    ("3️⃣ Ajouter les plugins", "npm install @capacitor/nfc @capacitor/push-notifications @capacitor/camera", "5 min", "Ce sont les \"traducteurs\" entre le web et le téléphone. NFC = lire les tags sans contact."),
    ("4️⃣ Adapter le code NFC", "Ajouter ~30 lignes de JS pour lire les tags NFC", "2h", "Le code le plus important : quand on tape un téléphone, lire l'ID et valider la mission."),
    ("5️⃣ Build + Sync", "npm run build → npx cap copy → npx cap sync", "10 min", "Copie ton site Next.js buildé dans les dossiers iOS/Android."),
    ("6️⃣ Tester sur téléphone", "npx cap open android → lancer sur ton vrai téléphone", "1h", "Tu branches ton téléphone en USB, tu lances l'app dessus."),
    ("7️⃣ Générer l'APK", "Android Studio → Build → Generate Signed Bundle", "1h", "Le fichier .apk ou .aab à envoyer aux testeurs ou à publier."),
    ("8️⃣ Compte développeur Apple", "inscription sur developer.apple.com (99€/an)", "30 min", "Obligatoire pour installer sur iPhone. Android n'a pas ce coût (25$ unique)."),
]
for i, (title, cmd, time, desc) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.7) + Inches(row * 1.35)
    card(s, x, y, Inches(5.8), Inches(1.2))
    # Accent bar
    rect(s, x, y, Inches(5.8), Inches(0.04), ACCENT)
    tf = tb(s, x + Inches(0.2), y + Inches(0.1), Inches(4), Inches(0.3))
    P(tf, title, 13, True, TEXT)
    tf = tb(s, x + Inches(4.2), y + Inches(0.1), Inches(1.4), Inches(0.3))
    P(tf, f"⏱️ {time}", 10, True, ACCENT, PP_ALIGN.RIGHT)
    # Command in monospace style
    rect(s, x + Inches(0.2), y + Inches(0.4), Inches(5.4), Inches(0.3), SURFACE2, None, 0.01)
    tf = tb(s, x + Inches(0.3), y + Inches(0.4), Inches(5.2), Inches(0.3))
    P(tf, cmd, 8, False, TEXT_SEC, name="Consolas")
    tf = tb(s, x + Inches(0.2), y + Inches(0.75), Inches(5.4), Inches(0.4))
    P(tf, desc, 9, False, TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 - REACT NATIVE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, AMBER)
tf = tb(s, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7))
P(tf, "⚛️ Option 3 : React Native (avec Expo)", 32, True, WHITE, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(12), Inches(1))
P(tf, "C'est quoi ?", 16, True, RGBColor(0xFE, 0xF3, 0xC7))
P(tf, "Contrairement aux deux options précédentes, React Native ne met PAS ton site web dans une coquille.", 12, False, RGBColor(0xFD, 0xF0, 0xB0))
P(tf, "Il REconstruit toute l'interface en composants 100% natifs. C'est comme si tu réécrivais TimeHeroes en \"langage iPhone/Android\" plutôt qu'en HTML.", 12, False, RGBColor(0xFD, 0xF0, 0xB0))

# What stays / what changes
rect(s, Inches(0.8), Inches(3), Inches(5.8), Inches(3.5), RGBColor(0xFD, 0xF0, 0xB0), RGBColor(0xE8, 0xC0, 0x50), 0.02)
tf = tb(s, Inches(1.1), Inches(3.1), Inches(5.2), Inches(3.3))
P(tf, "✅ Ce qui RESTE inchangé :", 16, True, WHITE, name="Calibri Light")
P(tf, "• Le serveur Next.js (API, base de données, Prisma)", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "• Toutes les routes /api/services, /api/bookings…", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "• La base SQLite avec toutes les données", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "", 8)
P(tf, "❌ Ce qui CHANGE :", 16, True, WHITE, name="Calibri Light")
P(tf, "• L'interface utilisateur (tout l'écran)", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "• Les composants (Dashboard, Services, Bookings)", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "• La navigation (plus de pages web, des \"écrans\")", 12, False, RGBColor(0xFF, 0xF5, 0xD0))

rect(s, Inches(7), Inches(3), Inches(5.8), Inches(3.5), RGBColor(0xFD, 0xF0, 0xB0), RGBColor(0xE8, 0xC0, 0x50), 0.02)
tf = tb(s, Inches(7.3), Inches(3.1), Inches(5.2), Inches(3.3))
P(tf, "💰 Budget détaillé si Ronald fait tout :", 16, True, WHITE, name="Calibri Light")
P(tf, "", 8)
P(tf, "• Apprentissage React Native : 1 semaine", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "  (Expo + React Navigation + composants de base)", 10, False, RGBColor(0xD0, 0xC0, 0x80))
P(tf, "• Réécriture des écrans : 1-2 semaines", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "  (Dashboard, Missions, Booking, Profil, etc.)", 10, False, RGBColor(0xD0, 0xC0, 0x80))
P(tf, "• Test + corrections : 1 semaine", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "• Publication App Store + Play Store : 2 jours", 12, False, RGBColor(0xFF, 0xF5, 0xD0))
P(tf, "", 6)
P(tf, "💰 TOTAL : 2 à 4 semaines • 99€ Apple + 25$ Google", 13, True, WHITE)

rect(s, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.5), RGBColor(0xEC, 0xFD, 0xF5), ACCENT, 0.02)
tf = tb(s, Inches(1), Inches(6.8), Inches(11), Inches(0.5))
P(tf, "⏱️ Estimation si Ronald fait tout : 2-4 semaines  •  💰 Coûts : 99€ + 25$ + 0€ de main-d'œuvre  •  📱 App 100% native", 12, True, RGBColor(0x85, 0x64, 0x04), PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 - TABLEAU COMPARATIF
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "📊 Comparatif — si Ronald fait tout lui-même", 30, True, TEXT, name="Calibri Light")

headers = ["", "PWA", "Capacitor", "React Native"]
col_x = [Inches(0.8), Inches(4), Inches(6.8), Inches(9.6)]
col_w = [Inches(3), Inches(2.6), Inches(2.6), Inches(2.6)]
for i, h in enumerate(headers):
    rect(s, col_x[i], Inches(1.2), col_w[i], Inches(0.45), ACCENT_DARK)
    tf = tb(s, col_x[i] + Inches(0.1), Inches(1.2), col_w[i] - Inches(0.2), Inches(0.45))
    P(tf, h, 11, True, WHITE, PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

rows = [
    ("Temps (Ronald seul)", "1 jour ⚡", "3-4 jours", "2-4 semaines"),
    ("Code à réécrire", "0%", "~5% (plugins)", "~60% (toute l'UI)"),
    ("NFC ?", "❌", "✅", "✅"),
    ("Notifications push", "✅ Web", "✅ Natives", "✅ Natives"),
    ("Hors-ligne", "✅", "✅", "✅"),
    ("Compte Apple (99€/an)", "❌", "✅ OBLIGATOIRE", "✅ OBLIGATOIRE"),
    ("Compte Google (25$)", "❌", "✅", "✅"),
    ("App Store / Play Store", "❌ Pas besoin", "✅ Publiable", "✅ Publiable"),
    ("Difficulté technique", "Facile", "Moyen", "Avancé"),
    ("Quand choisir ?", "Démo rapide", "🚀 RECOMMANDÉ", "Scale 10k+ users"),
]
for ri, (label, *vals) in enumerate(rows):
    y = Inches(1.7) + Inches(ri * 0.48)
    bg = SURFACE if ri % 2 == 0 else SURFACE2
    is_rec = "RECOMMANDÉ" in label or "RECOMMANDÉ" in vals[0] or "RECOMMANDÉ" in vals[1] or "RECOMMANDÉ" in vals[2]
    rect(s, Inches(0.8), y, Inches(11.7), Inches(0.45), bg, ACCENT if is_rec else BORDER if is_rec else None)
    tf = tb(s, Inches(0.9), y + Inches(0.05), Inches(3), Inches(0.35))
    P(tf, label, 10, True, TEXT)
    for vi, v in enumerate(vals):
        is_best = "⚡" in v or "RECOMMANDÉ" in v
        c = ACCENT_DARK if is_best else (GREEN_DARK if "✅" in v else (RGBColor(0xDC, 0x26, 0x26) if "❌" in v else TEXT_SEC))
        tf = tb(s, col_x[vi+1] + Inches(0.1), y + Inches(0.05), col_w[vi+1] - Inches(0.2), Inches(0.35))
        P(tf, v, 10, is_best or "RECOMMANDÉ" in v, c, PP_ALIGN.CENTER)

# Highlight
rect(s, Inches(0.8), Inches(6.38), Inches(11.7), Inches(0.55), RGBColor(0xEC, 0xFD, 0xF5), ACCENT, 0.02)
tf = tb(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.5))
P(tf, "🎯 Verdict : Capacitor est le meilleur rapport temps/bénéfice pour TimeHeroes. NFC natif + 3-4 jours de travail + 0€ de main-d'œuvre.", 13, True, ACCENT_DARK, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 - FOCUS NFC
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "📡 Focus : NFC — c'est quoi ?", 30, True, TEXT, name="Calibri Light")
tf = tb(s, Inches(0.8), Inches(1.1), Inches(12), Inches(0.5))
P(tf, "Le NFC (Near Field Communication) est la technologie derrière le paiement sans contact. Même principe, mais pour valider les missions.", 14, False, TEXT_SEC)

# How NFC works
rect(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), SURFACE, BORDER, 0.02)
tf = tb(s, Inches(1.1), Inches(1.9), Inches(5), Inches(4.8))
P(tf, "📱 Le scénario NFC pour TimeHeroes :", 16, True, TEXT, name="Calibri Light")
P(tf, "", 6)
P(tf, "1️⃣ Karim (bénévole) vient chez Mme Martin pour l'aider avec son smartphone.", 11, False, TEXT_SEC)
P(tf, "2️⃣ La mission est terminée. Karim ouvre l'app TimeHeroes → clique sur \"Valider\"", 11, False, TEXT_SEC)
P(tf, "3️⃣ Mme Martin approche son téléphone du téléphone de Karim (comme pour payer en sans contact)", 11, False, TEXT_SEC)
P(tf, "4️⃣ *BIP* — La mission est validée ! Les TIME sont libérés.", 11, True, ACCENT)
P(tf, "", 6)
P(tf, "Avantages du NFC par rapport au QR code :", 14, True, TEXT, name="Calibri Light")
P(tf, "• Plus rapide : 0,5 seconde au lieu de 5 secondes", 11, False, TEXT_SEC)
P(tf, "• Plus facile : pas besoin d'ouvrir l'app, d'allumer l'écran", 11, False, TEXT_SEC)
P(tf, "• Plus fiable : marche même en plein soleil (pas de reflet)", 11, False, TEXT_SEC)
P(tf, "• Idéal pour les seniors (pas de manipulation complexe)", 11, False, TEXT_SEC)

# QR comparison
rect(s, Inches(6.7), Inches(1.8), Inches(5.8), Inches(2.3), SURFACE, BORDER, 0.02)
tf = tb(s, Inches(7), Inches(1.9), Inches(5.2), Inches(2.1))
P(tf, "📷 QR Code (situation actuelle)", 14, True, ACCENT, name="Calibri Light")
P(tf, "• Ouvrir l'app SMS / WhatsApp", 10, False, TEXT_SEC)
P(tf, "• Le prestataire génère un QR code sur son écran", 10, False, TEXT_SEC)
P(tf, "• Le bénéficiaire scanne avec son appareil photo", 10, False, TEXT_SEC)
P(tf, "• Cliquer sur le lien → validation", 10, False, TEXT_SEC)
P(tf, "• ⏱️ ~5-10 secondes", 10, True, AMBER)

rect(s, Inches(6.7), Inches(4.3), Inches(5.8), Inches(2.5), SURFACE, BORDER, 0.02)
tf = tb(s, Inches(7), Inches(4.4), Inches(5.2), Inches(2.3))
P(tf, "📡 NFC (avec Capacitor)", 14, True, PURPLE, name="Calibri Light")
P(tf, "• Ouvrir l'app TimeHeroes", 10, False, TEXT_SEC)
P(tf, "• Cliquer sur \"Valider par NFC\"", 10, False, TEXT_SEC)
P(tf, "• Taper les deux téléphones l'un contre l'autre", 10, False, TEXT_SEC)
P(tf, "• *BIP*  Validé !", 10, True, ACCENT)
P(tf, "• ⏱️ ~1 seconde", 10, True, ACCENT)
P(tf, "", 4)
P(tf, "💰 Coût : ~0,30€ par tag NFC autocollant", 11, False, TEXT_SEC)
P(tf, "   (on peut coller le tag derrière le téléphone)", 10, False, TEXT_MUTED)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 - ROADMAP
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "🗺️ Roadmap — si Ronald fait tout", 30, True, TEXT, name="Calibri Light")

phases = [
    ("Phase 1", "PWA", "Jour 1", "0€", ACCENT, [
        "Ajouter manifest.json et les icônes",
        "Installer le service worker (@serwist/next)",
        "✅ TimeHeroes s'installe sur le téléphone",
        "✅ 1 jour de travail, zéro euro",
    ]),
    ("Phase 2", "Capacitor", "Jours 2-5", "~124€", PURPLE, [
        "Installer Capacitor et les plugins NFC/Push",
        "Ajouter ~30 lignes de code pour le NFC",
        "Tester la validation par tap sur mon téléphone",
        "Générer l'APK Android (publiable)",
        "✅ App complète avec NFC !",
    ]),
    ("Phase 3", "React Native", "Sem. 3-6", "~124€", AMBER, [
        "Apprendre React Native / Expo",
        "Réécrire les écrans un par un",
        "Tester sur iOS et Android",
        "Publier sur l'App Store et Google Play",
        "✅ App 100% native, premium",
    ]),
]
for i, (phase, name, time, cost, accent_color, steps) in enumerate(phases):
    x = Inches(0.8) + Inches(i * 4.2)
    card(s, x, Inches(1.3), Inches(3.8), Inches(5.8))
    rect(s, x, Inches(1.3), Inches(3.8), Inches(0.06), accent_color)
    tf = tb(s, x + Inches(0.2), Inches(1.5), Inches(3.4), Inches(0.3))
    P(tf, f"{phase}", 10, True, accent_color)
    tf = tb(s, x + Inches(0.2), Inches(1.8), Inches(3.4), Inches(0.4))
    P(tf, f"{name}", 18, True, TEXT, name="Calibri Light")
    rect(s, x + Inches(0.2), Inches(2.3), Inches(1.5), Inches(0.3), accent_color)
    tf = tb(s, x + Inches(0.25), Inches(2.3), Inches(1.4), Inches(0.3))
    P(tf, f"⏱️ {time}", 8, True, WHITE, PP_ALIGN.CENTER)
    rect(s, x + Inches(1.8), Inches(2.3), Inches(1.2), Inches(0.3), RGBColor(0xEC, 0xFD, 0xF5))
    tf = tb(s, x + Inches(1.85), Inches(2.3), Inches(1.1), Inches(0.3))
    P(tf, f"{cost}", 8, True, ACCENT_DARK, PP_ALIGN.CENTER)
    for j, step in enumerate(steps):
        icon = "  " if "✅" in step else "  "
        tf = tb(s, x + Inches(0.2), Inches(2.9) + Inches(j * 0.55), Inches(3.4), Inches(0.5))
        P(tf, f"{icon}  {step.replace('✅ ', '')}", 10, False, TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 - GLOSSAIRE SIMPLIFIÉ
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
P(tf, "📖 Glossaire — les mots techniques expliqués simplement", 30, True, TEXT, name="Calibri Light")

glossary = [
    ("PWA", "Progressive Web App", "C'est un site web normal, mais le navigateur propose de l'\"installer\" sur l'écran d'accueil. Ça a une icône, ça peut envoyer des notifications et marcher un peu hors-ligne. Exemple : Twitter, Pinterest."),
    ("Service Worker", "", "Un petit programme qui tourne en arrière-plan sur le téléphone, même quand le site est fermé. Il gère le cache (hors-ligne) et les notifications push."),
    ("WebView", "", "C'est un navigateur Internet (Chrome/Safari) embarqué à l'intérieur d'une app. Quand tu ouvres un lien dans Instagram ou WhatsApp, c'est une WebView. Capacitor utilise une WebView pour afficher ton site."),
    ("Capacitor", "", "Un outil qui prend ton site web et le transforme en vraie app iPhone/Android. Il ajoute une \"coquille\" autour et permet d'utiliser le NFC, la caméra, les notifications natives."),
    ("React Native", "", "Un framework (créé par Meta/Facebook) pour construire des apps mobiles 100% natives. Contrairement à Capacitor, il ne met PAS ton site dans une WebView : il reconstruit toute l'interface en langage natif."),
    ("Expo", "", "Une version simplifiée de React Native. Ça permet de tester l'app directement sur son téléphone sans installation compliquée, et de builder dans le cloud (pas besoin d'avoir un Mac pour iOS)."),
    ("NFC", "", "Near Field Communication. Technologie sans contact utilisée pour le paiement par carte bancaire ou Apple Pay/Google Pay. Distance max : 4 cm. Permet de valider une mission en tapant deux téléphones."),
    ("Plugin", "", "Un \"module\" qu'on ajoute à Capacitor ou React Native pour accéder à une fonction du téléphone. Exemple : plugin NFC, plugin Caméra, plugin Notifications."),
    ("Tag NFC", "", "Une petite étiquette (type autocollant, ~0,30€) qui contient une puce NFC. On peut la coller derrière un téléphone. Quand un autre téléphone la tape, il lit l'information stockée."),
    ("App Store / Play Store", "", "Les magasins d'applications. App Store = iPhone (Apple, 99€/an pour publier). Google Play = Android (25$ unique). Sans ces stores, tu peux quand même installer l'app manuellement (fichier .apk)."),
    ("Build", "", "L'étape qui transforme ton code en une app installable. Pour Next.js : c'est \"npm run build\". Pour Android : ça génère un .apk. Pour iOS : une archive .ipa."),
    ("Escrow", "", "Un système de séquestre : les TIME sont bloqués quand tu réserves, et libérés seulement quand la mission est validée. Ça protège tout le monde."),
]
for i, (term, eng, defn) in enumerate(glossary):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.2) + Inches(row * 1.05)
    card(s, x, y, Inches(5.8), Inches(0.95))
    tf = tb(s, x + Inches(0.15), y + Inches(0.05), Inches(2.5), Inches(0.3))
    P(tf, term, 14, True, ACCENT)
    if eng:
        P(tf, f"= {eng}", 10, False, TEXT_MUTED)
    tf = tb(s, x + Inches(0.15), y + Inches(0.4), Inches(5.5), Inches(0.5))
    P(tf, defn, 9, False, TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 - CONCLUSION
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, ACCENT_DARK)
circle(s, Inches(-0.5), Inches(-0.5), Inches(2.5), RGBColor(0x00, 0xAA, 0x88))
circle(s, Inches(11), Inches(5), Inches(3.5), RGBColor(0x00, 0x99, 0x77))

tf = tb(s, Inches(1.5), Inches(1), Inches(10), Inches(1))
P(tf, "🎯 Conclusion : par quoi commencer ?", 40, True, WHITE, PP_ALIGN.LEFT, "Calibri Light")

tf = tb(s, Inches(1.5), Inches(2.2), Inches(10), Inches(4.5))
P(tf, "Semaine 1 — PWA (1 jour)", 24, True, ACCENT, name="Calibri Light")
P(tf, "Tu ajoutes manifest.json + service worker à Next.js.", 13, False, RGBColor(0xBB, 0xF7, 0xE0))
P(tf, "TimeHeroes devient installable sur le téléphone.", 13, False, RGBColor(0xBB, 0xF7, 0xE0))
P(tf, "Idéal pour la démo fondation : \"Installez l'app, testez-la\"", 13, False, RGBColor(0xBB, 0xF7, 0xE0))
P(tf, "💰 0€, pas de compte Apple nécessaire.", 13, True, WHITE)
P(tf, "", 8)

P(tf, "Semaine 1 aussi — Capacitor (3-4 jours)", 24, True, AMBER, name="Calibri Light")
P(tf, "Tu ajoutes Capacitor et les plugins NFC.", 13, False, RGBColor(0xBB, 0xF7, 0xE0))
P(tf, "La validation par tap NFC devient possible — argument DIFFÉRENCIANT pour la fondation.", 13, False, RGBColor(0xBB, 0xF7, 0xE0))
P(tf, "Tu gardes 95% de ton code. Tu apprends Capacitor en 1 jour.", 13, False, RGBColor(0xBB, 0xF7, 0xE0))
P(tf, "💰 99€ pour le compte Apple + 0,30€/tag NFC.", 13, True, WHITE)
P(tf, "", 8)

P(tf, "Plus tard — React Native (si 10 000+ utilisateurs)", 22, True, RGBColor(0xAA, 0xAA, 0xAA), name="Calibri Light")
P(tf, "Pour le moment, PAS nécessaire. Capacitor suffit largement.", 12, False, RGBColor(0xAA, 0xAA, 0xAA))

tf = tb(s, Inches(1.5), Inches(6.5), Inches(10), Inches(0.5))
P(tf, "Prêt pour la Phase 1 ?  🚀", 18, True, WHITE)


# ═══════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════
out = "/root/projects/timebank-poc/Cours-App-Mobile-TimeHeroes-v2.pptx"
prs.save(out)
print(f"✅ Saved: {out}")
print(f"📊 {len(prs.slides)} slides")
