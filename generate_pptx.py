"""TimeHeroes — Cours : Transformer le site en app mobile (powerpoint)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── PALETTE LIGHT (TimeHeroes) ──────────────────────────────────────
BG = RGBColor(0xF8, 0xF5, 0xF0)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2 = RGBColor(0xF0, 0xED, 0xE8)
ACCENT = RGBColor(0x2B, 0xB2, 0x86)
ACCENT_DARK = RGBColor(0x00, 0x8F, 0x78)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_SEC = RGBColor(0x6B, 0x6B, 0x6B)
TEXT_MUTED = RGBColor(0x9A, 0x9A, 0x9A)
BORDER = RGBColor(0xE5, 0xE0, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PINK = RGBColor(0xEC, 0x48, 0x99)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x10, 0xB9, 0x81)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]  # blank layout


# ─── HELPERS ─────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color=WHITE, border=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if radius:
        s.adjustments[0] = radius
    return s

def circle(slide, l, t, size, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def textbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb.text_frame

def p_add(tf, text="", size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, name="Calibri"):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '' and not text:
        p = tf.paragraphs[0]
    elif len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = name
    p.alignment = align
    p.space_after = Pt(4)
    return p

def add_run(p, text, size=14, bold=False, color=TEXT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r

def card(slide, l, t, w, h, title="", body="", accent=None):
    c = rect(slide, l, t, w, h, SURFACE, BORDER, 0.02)
    if accent:
        rect(slide, l + Inches(0.05), t + Inches(0.1), Inches(0.04), h - Inches(0.2), accent)
    tf = textbox(slide, l + Inches(0.25), t + Inches(0.15), w - Inches(0.5), h - Inches(0.3))
    if title:
        p_add(tf, title, 14, True, TEXT)
    if body:
        p_add(tf, body, 11, False, TEXT_SEC)


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG)
    return slide


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, ACCENT_DARK)
# Decorative circles
circle(s, Inches(-1), Inches(-1), Inches(3), RGBColor(0x22, 0xCC, 0xA0))
circle(s, Inches(11), Inches(5), Inches(4), RGBColor(0x00, 0xAA, 0x88))
circle(s, Inches(10), Inches(-2), Inches(2.5), RGBColor(0x00, 0x99, 0x77))

tf = textbox(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1))
p_add(tf, "📱 Transformer TimeHeroes en App Mobile", 40, True, WHITE, PP_ALIGN.LEFT, "Calibri Light")

tf = textbox(s, Inches(1.5), Inches(3), Inches(10), Inches(1.5))
p_add(tf, "Cours complet : les 3 options pour passer du web au mobile\nPWA · Capacitor · React Native", 18, False, RGBColor(0xCC, 0xFA, 0xE8), PP_ALIGN.LEFT)

# Infos
tf = textbox(s, Inches(1.5), Inches(5.5), Inches(5), Inches(1))
p_add(tf, "Projet : TimeHeroes — Banque du Temps", 12, False, RGBColor(0xAA, 0xEE, 0xDD))
p_add(tf, "Stack actuel : Next.js 16 · Prisma · SQLite · Caddy", 12, False, RGBColor(0xAA, 0xEE, 0xDD))
p_add(tf, "Juillet 2026", 12, False, RGBColor(0xAA, 0xEE, 0xDD))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — POURQUOI PASSER EN MOBILE ?
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
# Title
tf = textbox(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6))
p_add(tf, "🎯 Pourquoi passer en mobile ?", 30, True, TEXT)
# Sub
tf = textbox(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5))
p_add(tf, "Le web ça marche, mais le mobile ouvre des super-pouvoirs", 14, False, TEXT_SEC)

data = [
    ("📷 Scanner un QR code", "❌ Pas de caméra native", "✅ Scan natif"),
    ("📱 Valider par NFC", "❌ Pas de NFC en Web", "✅ Tap to validate"),
    ("🔔 Notifications push", "⚠️ Si onglet ouvert", "✅ Push natif toujours"),
    ("📴 Mode hors-ligne", "❌ Pas de cache durable", "✅ Cache offline"),
    ("🏠 Icône écran d'accueil", "⚠️ Dépend du navigateur", "✅ Permanente"),
]
for i, (label, web, mobile) in enumerate(data):
    y = Inches(2) + Inches(i * 1)
    card(s, Inches(0.8), y, Inches(5.5), Inches(0.85), "", label, None)
    rect(s, Inches(6.5), y, Inches(2.5), Inches(0.85), RGBColor(0xFE, 0xF0, 0xF0), RGBColor(0xFE, 0xCB, 0xCB))
    tf = textbox(s, Inches(6.6), y + Inches(0.15), Inches(2.3), Inches(0.6))
    p_add(tf, web, 11, False, RGBColor(0xB9, 0x1C, 0x1C), PP_ALIGN.CENTER)
    rect(s, Inches(9.2), y, Inches(2.5), Inches(0.85), RGBColor(0xEC, 0xFD, 0xF5), RGBColor(0x6E, 0xE7, 0xB7))
    tf = textbox(s, Inches(9.3), y + Inches(0.15), Inches(2.3), Inches(0.6))
    p_add(tf, mobile, 11, True, ACCENT_DARK, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — LES 3 CHEMINS
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = textbox(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6))
p_add(tf, "🛤️ Les 3 chemins possibles", 30, True, TEXT)
tf = textbox(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.4))
p_add(tf, "Trois stratégies complémentaires, pas concurrentes", 14, False, TEXT_SEC)

options = [
    ("PWA", "1 jour · 0€", "Site web qui s'installe\ncomme une app", ACCENT, "✅ Zéro coût\n✅ 100% code réutilisé\n❌ Pas de NFC"),
    ("Capacitor", "3-5 jours · ~124€/an", "WebView native\n+ pont NFC/Camera", PURPLE, "✅ NFC natif\n✅ 95% code réutilisé\n❌ Perf WebView"),
    ("React Native", "4-8 sem · 10k-32k€", "App 100% native\nréécrite de zéro", AMBER, "✅ Performance max\n✅ NFC + tout natif\n❌ 2 codebases"),
]
for i, (name, time, desc, accent_color, pros) in enumerate(options):
    x = Inches(0.8) + Inches(i * 4.2)
    # Card
    c = rect(s, x, Inches(1.9), Inches(3.8), Inches(5), SURFACE, BORDER, 0.02)
    # Accent top bar
    rect(s, x, Inches(1.9), Inches(3.8), Inches(0.06), accent_color)
    # Circle with number
    circle(s, x + Inches(1.4), Inches(2.2), Inches(1), accent_color)
    tf = textbox(s, x + Inches(1.4), Inches(2.2), Inches(1), Inches(1))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p_add(tf, str(i+1), 32, True, WHITE, PP_ALIGN.CENTER, "Calibri Light")
    # Name
    tf = textbox(s, x + Inches(0.3), Inches(3.4), Inches(3.2), Inches(0.5))
    p_add(tf, name, 22, True, TEXT, PP_ALIGN.CENTER, "Calibri Light")
    # Time
    tf = textbox(s, x + Inches(0.3), Inches(3.9), Inches(3.2), Inches(0.4))
    p_add(tf, time, 12, True, accent_color, PP_ALIGN.CENTER)
    # Description
    tf = textbox(s, x + Inches(0.3), Inches(4.3), Inches(3.2), Inches(0.6))
    p_add(tf, desc, 11, False, TEXT_SEC, PP_ALIGN.CENTER)
    # Pros/cons
    tf = textbox(s, x + Inches(0.3), Inches(5.0), Inches(3.2), Inches(1.5))
    for line in pros.split("\n"):
        p_add(tf, line, 10, False, TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — PWA DÉTAIL
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
# Decorative bg
rect(s, Inches(0), Inches(0), Inches(4.5), H, ACCENT)
tf = textbox(s, Inches(0.8), Inches(0.6), Inches(3.2), Inches(1.5))
p_add(tf, "📲", 48, False, WHITE)
p_add(tf, "PWA", 36, True, WHITE, name="Calibri Light")
p_add(tf, "Progressive Web App", 14, False, RGBColor(0xBB, 0xF7, 0xE0))
p_add(tf, "1 jour · 0€", 12, True, RGBColor(0x88, 0xF0, 0xD0))

# Content on the right
x = Inches(5)
tf = textbox(s, x, Inches(0.5), Inches(8), Inches(0.5))
p_add(tf, "Un site web qui se comporte comme une app", 22, True, TEXT, name="Calibri Light")

items = [
    "📄 manifest.json → icône + thème + écran d'accueil",
    "⚙️ Service Worker → cache offline + notifications push",
    "📱 ✅ Fonctionne hors-ligne (pages en cache)",
    "🔔 ✅ Notifications push (Web Push API)",
    "📷 ⚠️ QR code via caméra Web (getUserMedia)",
    "❌ NFC : impossible en WebView",
]
for i, item in enumerate(items):
    y = Inches(1.3) + Inches(i * 0.65)
    rect(s, x, y, Inches(7.8), Inches(0.5), SURFACE, BORDER, 0.02)
    tf = textbox(s, x + Inches(0.2), y + Inches(0.08), Inches(7.4), Inches(0.35))
    p_add(tf, item, 12, False, TEXT)

# Examples
tf = textbox(s, x, Inches(5.5), Inches(8), Inches(0.4))
p_add(tf, "🌍 Exemples célèbres de PWA", 16, True, TEXT, name="Calibri Light")
ex = ["Twitter (X)", "Pinterest (+40% de temps)", "Starbucks (commande offline)", "Uber (version allégée)"]
for i, e in enumerate(ex):
    tf = textbox(s, x + Inches(i * 2), Inches(6), Inches(1.9), Inches(0.5))
    rect(s, x + Inches(i * 2), Inches(6), Inches(1.9), Inches(0.5), SURFACE, BORDER, 0.02)
    tf = textbox(s, x + Inches(i * 2) + Inches(0.1), Inches(6) + Inches(0.08), Inches(1.7), Inches(0.35))
    p_add(tf, e, 10, False, TEXT_SEC, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — CAPACITOR DÉTAIL
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, Inches(0), Inches(0), Inches(4.5), H, PURPLE)
tf = textbox(s, Inches(0.8), Inches(0.6), Inches(3.2), Inches(1.5))
p_add(tf, "🔌", 48, False, WHITE)
p_add(tf, "Capacitor", 36, True, WHITE, name="Calibri Light")
p_add(tf, "Wrapper natif autour du site", 14, False, RGBColor(0xDD, 0xD6, 0xFC))
p_add(tf, "3-5 jours · ~124€/an", 12, True, RGBColor(0xCC, 0xC0, 0xF5))

x = Inches(5)
tf = textbox(s, x, Inches(0.5), Inches(8), Inches(0.5))
p_add(tf, "Le pont entre le web et le natif", 22, True, TEXT, name="Calibri Light")

# Architecture diagram
rect(s, x, Inches(1.2), Inches(7.8), Inches(2.5), SURFACE, BORDER, 0.02)
items_y = [
    (Inches(0.3), "🌐 WebView (ton site Next.js)", ACCENT),
    (Inches(0.9), "⬇️  Pont JS → Natif (Capacitor Plugins)", PURPLE),
    (Inches(1.5), "📱 Android (Kotlin)  ·  iOS (Swift)", AMBER),
]
for dy, label, color in items_y:
    tf = textbox(s, x + Inches(0.3), Inches(1.3) + dy, Inches(7), Inches(0.5))
    p_add(tf, label, 13, True if "⬇️" in label else False, color if "⬇️" in label else TEXT)

# Feature grid
features = [
    ("📱 NFC natif", "✅ Validation par tap", ACCENT),
    ("🔔 Push natif", "✅ FCM + APN", ACCENT),
    ("📷 Caméra QR", "✅ Plugin Capacitor", ACCENT),
    ("📴 Offline", "✅ Workbox cache", ACCENT),
    ("💾 95% code réutilisé", "✅ Next.js inchangé", ACCENT),
    ("🔄 Mise à jour", "✅ rebuild + npx cap sync", ACCENT),
]
for i, (feat, desc, color) in enumerate(features):
    fx = x + Inches(i % 3 * 2.7)
    fy = Inches(4.2) + Inches(i // 3 * 1.2)
    rect(s, fx, fy, Inches(2.5), Inches(1), SURFACE, BORDER, 0.02)
    tf = textbox(s, fx + Inches(0.15), fy + Inches(0.1), Inches(2.2), Inches(0.8))
    p_add(tf, feat, 12, True, color)
    p_add(tf, desc, 10, False, TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — REACT NATIVE DÉTAIL
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
rect(s, Inches(0), Inches(0), Inches(4.5), H, AMBER)
tf = textbox(s, Inches(0.8), Inches(0.6), Inches(3.2), Inches(1.5))
p_add(tf, "⚛️", 48, False, WHITE)
p_add(tf, "React Native", 36, True, WHITE, name="Calibri Light")
p_add(tf, "App 100% native (réécriture UI)", 14, False, RGBColor(0xFE, 0xF3, 0xC7))
p_add(tf, "4-8 semaines · 10k-32k€", 12, True, RGBColor(0xFD, 0xE6, 0x8A))

x = Inches(5)
tf = textbox(s, x, Inches(0.5), Inches(8), Inches(0.5))
p_add(tf, "Le grand frère premium", 22, True, TEXT, name="Calibri Light")

# Architecture
rect(s, x, Inches(1.2), Inches(7.8), Inches(2), SURFACE, BORDER, 0.02)
tf = textbox(s, x + Inches(0.3), Inches(1.35), Inches(7.2), Inches(1.7))
p_add(tf, "Architecture :", 14, True, TEXT)
p_add(tf, "┌──────────────────────────────┐", 10, False, TEXT_MUTED)
p_add(tf, "│  App React Native (Expo)       │  ← UI réécrite", 10, False, TEXT)
p_add(tf, "│  ├── Écrans natifs (FlatList…)  │", 10, False, TEXT)
p_add(tf, "│  ├── NFC · Camera · Push        │", 10, False, TEXT)
p_add(tf, "│  └── Appelle /api/ de Next.js   │", 10, False, TEXT)
p_add(tf, "└──────────────────────────────┘", 10, False, TEXT_MUTED)

# Pros/cons
tf = textbox(s, x, Inches(3.6), Inches(3.8), Inches(0.4))
p_add(tf, "✅ Avantages", 16, True, GREEN, name="Calibri Light")
pros = ["Performance native (60fps)", "UX mobile authentique", "Tout l'accès natif", "Expo EAS Build (cloud)"]
for i, p in enumerate(pros):
    tf = textbox(s, x + Inches(i % 2 * 2), Inches(4.1) + Inches(i // 2 * 0.5), Inches(1.9), Inches(0.45))
    rect(s, x + Inches(i % 2 * 2), Inches(4.1) + Inches(i // 2 * 0.5), Inches(1.9), Inches(0.45), SURFACE, BORDER, 0.02)
    tf = textbox(s, x + Inches(i % 2 * 2) + Inches(0.1), Inches(4.1) + Inches(i // 2 * 0.5) + Inches(0.08), Inches(1.7), Inches(0.3))
    p_add(tf, p, 10, False, TEXT)

tf = textbox(s, x + Inches(4), Inches(3.6), Inches(3.8), Inches(0.4))
p_add(tf, "❌ Inconvénients", 16, True, RGBColor(0xDC, 0x26, 0x26), name="Calibri Light")
cons = ["Temps : 4-8 semaines", "Deux codebases à maintenir", "Coût dev : 10k-32k€", "Fragmentation iOS/Android"]
for i, c in enumerate(cons):
    tf = textbox(s, x + Inches(4) + Inches(i % 2 * 2), Inches(4.1) + Inches(i // 2 * 0.5), Inches(1.9), Inches(0.45))
    rect(s, x + Inches(4) + Inches(i % 2 * 2), Inches(4.1) + Inches(i // 2 * 0.5), Inches(1.9), Inches(0.45), SURFACE, BORDER, 0.02)
    tf = textbox(s, x + Inches(4) + Inches(i % 2 * 2) + Inches(0.1), Inches(4.1) + Inches(i // 2 * 0.5) + Inches(0.08), Inches(1.7), Inches(0.3))
    p_add(tf, c, 10, False, TEXT)

# Famous apps
tf = textbox(s, x, Inches(5.5), Inches(8), Inches(0.4))
p_add(tf, "🏆 Apps célèbres en React Native", 16, True, TEXT, name="Calibri Light")
apps = ["Instagram", "Uber Eats", "Coinbase", "Walmart", "Discord", "Airbnb (avant)"]
for i, app in enumerate(apps):
    tf = textbox(s, x + Inches(i * 1.3), Inches(6), Inches(1.2), Inches(0.5))
    rect(s, x + Inches(i * 1.3), Inches(6), Inches(1.2), Inches(0.5), SURFACE, BORDER, 0.02)
    tf = textbox(s, x + Inches(i * 1.3) + Inches(0.05), Inches(6) + Inches(0.08), Inches(1.1), Inches(0.35))
    p_add(tf, app, 10, False, TEXT_SEC, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — TABLEAU COMPARATIF
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = textbox(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
p_add(tf, "📊 Comparatif détaillé", 30, True, TEXT, name="Calibri Light")

headers = ["Critère", "PWA", "Capacitor", "React Native"]
cols_x = [Inches(0.8), Inches(4), Inches(6.8), Inches(9.6)]
col_w = [Inches(3), Inches(2.6), Inches(2.6), Inches(2.6)]

# Header row
for i, h in enumerate(headers):
    rect(s, cols_x[i], Inches(1.2), col_w[i], Inches(0.45), ACCENT_DARK)
    tf = textbox(s, cols_x[i] + Inches(0.1), Inches(1.2), col_w[i] - Inches(0.2), Inches(0.45))
    p_add(tf, h, 11, True, WHITE, PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

rows = [
    ("Effort", "1 jour", "3-5 jours", "4-8 semaines"),
    ("Code réutilisé", "100%", "~95%", "~40% (API)"),
    ("NFC natif", "❌", "✅", "✅"),
    ("Notifications push", "✅ Web Push", "✅ FCM/APN", "✅ FCM/APN"),
    ("Performance", "Moyenne", "Bonne", "Excellente"),
    ("App Store", "❌", "✅", "✅"),
    ("Coût fixe", "0€", "~124€/an", "~124€/an + dev"),
    ("Coût dev (est.)", "0€ (toi)", "0€ (toi)", "10k-32k€"),
    ("Maintenance", "Faible", "Faible", "Élevée"),
    ("Recommandation", "POC/MVP", "✅ Prod NFC", "Scale 10k+"),
]
for ri, (label, *vals) in enumerate(rows):
    y = Inches(1.7) + Inches(ri * 0.48)
    bg = SURFACE if ri % 2 == 0 else SURFACE2
    rect(s, Inches(0.8), y, Inches(11.7), Inches(0.45), bg)
    tf = textbox(s, Inches(0.9), y + Inches(0.05), Inches(3), Inches(0.35))
    p_add(tf, label, 10, True, TEXT)
    for vi, v in enumerate(vals):
        tf = textbox(s, cols_x[vi+1] + Inches(0.1), y + Inches(0.05), col_w[vi+1] - Inches(0.2), Inches(0.35))
        color = ACCENT if v == "✅" or v == "✅ Prod NFC" else (RGBColor(0xDC, 0x26, 0x26) if v == "❌" else TEXT_SEC)
        p_add(tf, v, 10, "✅" in v, color, PP_ALIGN.CENTER)

# Highlight recommendation row
rect(s, Inches(0.8), Inches(6.38), Inches(11.7), Inches(0.45), RGBColor(0xEC, 0xFD, 0xF5), ACCENT)
tf = textbox(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.4))
p_add(tf, "🎯 Recommandé pour TimeHeroes : Capacitor (NFC natif, 95% de code conservé, déploiement rapide)", 12, True, ACCENT_DARK, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — NFC FOCUS
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, RGBColor(0x05, 0x2E, 0x16))
circle(s, Inches(-0.5), Inches(5), Inches(3), RGBColor(0x00, 0x66, 0x44))
circle(s, Inches(11), Inches(-1), Inches(2.5), RGBColor(0x00, 0x55, 0x33))

tf = textbox(s, Inches(0.8), Inches(0.5), Inches(12), Inches(0.6))
p_add(tf, "📡 Focus : NFC et QR Code", 30, True, WHITE, name="Calibri Light")
tf = textbox(s, Inches(0.8), Inches(1.2), Inches(12), Inches(0.4))
p_add(tf, "La validation des missions par tap NFC, un argument différenciant fort pour la fondation", 14, False, RGBColor(0xAA, 0xEE, 0xDD))

# QR vs NFC comparison
for i, (title, items) in enumerate([
    ("📷 QR Code (actuel)", [
        "📸 Scan avec la caméra",
        "⏱️ ~5 secondes par scan",
        "💰 0€ (écran uniquement)",
        "📱 Tous les téléphones",
        "⚠️ Moins confort pour les seniors",
    ]),
    ("📡 NFC (futur)", [
        "🤝 Tap à ~4 cm",
        "⏱️ ~0.5 seconde",
        "💰 ~0,30€ par tag NFC",
        "📱 Smartphones récents",
        "✅ Marche même téléphone éteint",
    ]),
]):
    x = Inches(1) + Inches(i * 6)
    rect(s, x, Inches(1.9), Inches(5.5), Inches(3.5), RGBColor(0x0A, 0x0A, 0x0A), RGBColor(0x15, 0x15, 0x15), 0.02)
    tf = textbox(s, x + Inches(0.3), Inches(2.1), Inches(4.9), Inches(0.5))
    p_add(tf, title, 18, True, ACCENT if i == 1 else WHITE, name="Calibri Light")
    for j, item in enumerate(items):
        tf = textbox(s, x + Inches(0.3), Inches(2.8) + Inches(j * 0.45), Inches(4.9), Inches(0.4))
        p_add(tf, item, 11, False, RGBColor(0xCC, 0xCC, 0xCC))

# Recommendation
rect(s, Inches(1), Inches(5.8), Inches(11.3), Inches(1.2), RGBColor(0x0A, 0x2E, 0x1A), ACCENT, 0.02)
tf = textbox(s, Inches(1.3), Inches(5.9), Inches(10.7), Inches(1))
p_add(tf, "💡 Recommandation : Garder les deux", 16, True, WHITE)
p_add(tf, "QR code pour le POC immédiat, NFC en upgrade pour la version Capacitor.", 12, False, RGBColor(0xBB, 0xF7, 0xE0))
p_add(tf, "Les tags NFC coûtent 0,30€ pièce → 30€ pour 100 tags. Les seniors n'auront qu'à taper leur téléphone.", 12, False, RGBColor(0xBB, 0xF7, 0xE0))


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = textbox(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
p_add(tf, "🗺️ Roadmap recommandée pour TimeHeroes", 30, True, TEXT, name="Calibri Light")

phases = [
    ("Phase 1", "PWA", "1 jour · 0€", "Démo fondation", ACCENT, [
        "Ajouter manifest.json à Next.js",
        "Installer @serwist/next (service worker)",
        "Ajouter les icônes 192×192 et 512×512",
        "✅ Ajouter à l'écran d'accueil",
    ]),
    ("Phase 2", "Capacitor", "3-5 jours · 124€", "Pilote CCAS", PURPLE, [
        "Initialiser Capacitor dans le projet",
        "Plugin NFC + Push + Camera",
        "Adapter validation QR → NFC + QR",
        "Build APK + Archive iOS",
        "✅ App complète NFC",
    ]),
    ("Phase 3", "React Native", "4-8 sem · 10k€+", "Scale 10k+ users", AMBER, [
        "Nouveau projet Expo",
        "Réécrire les écrans clés",
        "Brancher sur l'API existante",
        "Animer avec Reanimated 2",
        "✅ App 100% native premium",
    ]),
]
for i, (phase, name, time, goal, accent_color, steps) in enumerate(phases):
    x = Inches(0.8) + Inches(i * 4.2)
    # Card
    rect(s, x, Inches(1.3), Inches(3.8), Inches(5.8), SURFACE, BORDER, 0.02)
    # Accent bar
    rect(s, x, Inches(1.3), Inches(3.8), Inches(0.06), accent_color)
    # Phase header
    tf = textbox(s, x + Inches(0.2), Inches(1.5), Inches(3.4), Inches(0.4))
    p_add(tf, f"{phase}", 11, True, accent_color)
    tf = textbox(s, x + Inches(0.2), Inches(1.9), Inches(3.4), Inches(0.4))
    p_add(tf, f"{name}  ·  {time}", 14, True, TEXT, name="Calibri Light")
    # Goal badge
    rect(s, x + Inches(0.2), Inches(2.4), Inches(2), Inches(0.35), accent_color)
    tf = textbox(s, x + Inches(0.25), Inches(2.4), Inches(1.9), Inches(0.35))
    p_add(tf, f"🎯 {goal}", 9, True, WHITE, PP_ALIGN.CENTER)
    # Steps
    for j, step in enumerate(steps):
        tf = textbox(s, x + Inches(0.2), Inches(3.0) + Inches(j * 0.5), Inches(3.4), Inches(0.45))
        p_add(tf, f"  {step}", 10, False, TEXT_SEC)

# Decision tree at bottom
tf = textbox(s, Inches(0.8), Inches(7.2), Inches(11.5), Inches(0.4))
p_add(tf, "💡  Arbre de décision : NFC nécessaire ? → OUI → Capacitor.  |  Budget ? → 0€ → PWA.  |  Scale 10k+ → React Native.", 11, False, TEXT_SEC, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — GLOSSAIRE
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = textbox(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
p_add(tf, "📖 Glossaire technique", 30, True, TEXT, name="Calibri Light")

terms = [
    ("PWA", "Progressive Web App — site web installable comme une app"),
    ("Service Worker", "Script JS en arrière-plan pour cache et push"),
    ("WebView", "Navigateur embarqué dans une app native"),
    ("Capacitor", "Pont web → natif (NFC, Camera, Push)"),
    ("React Native", "Framework pour apps natives en JavaScript/React"),
    ("NFC", "Near Field Communication — sans contact ~4 cm"),
    ("FCM / APN", "Firebase Cloud Msg / Apple Push Notification"),
    ("Escrow", "Séquestre : TIME bloqués jusqu'à validation"),
    ("Expo / EAS", "Surcouche React Native + build cloud"),
    ("Workbox", "Librairie Google pour la gestion du cache offline"),
]
for i, (term, defn) in enumerate(terms):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.3) + Inches(row * 1.1)
    rect(s, x, y, Inches(5.8), Inches(0.9), SURFACE, BORDER, 0.02)
    tf = textbox(s, x + Inches(0.2), y + Inches(0.08), Inches(5.4), Inches(0.3))
    p_add(tf, term, 12, True, ACCENT)
    tf = textbox(s, x + Inches(0.2), y + Inches(0.4), Inches(5.4), Inches(0.4))
    p_add(tf, defn, 10, False, TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — BUDGET
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
tf = textbox(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6))
p_add(tf, "💰 Budget estimé (année 1)", 30, True, TEXT, name="Calibri Light")

cols_x = [Inches(0.8), Inches(4.2), Inches(6.8), Inches(9.5)]
headers = ["Poste", "PWA", "Capacitor", "React Native"]
for i, h in enumerate(headers):
    rect(s, cols_x[i], Inches(1.3), Inches(3.2) if i == 0 else Inches(2.4), Inches(0.45), ACCENT_DARK)
    tf = textbox(s, cols_x[i] + Inches(0.1), Inches(1.3), (Inches(3.2) if i == 0 else Inches(2.4)) - Inches(0.2), Inches(0.45))
    p_add(tf, h, 11, True, WHITE, PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

rows_data = [
    ("Développement (toi)", "1j × 0€", "5j × 0€", "—"),
    ("Développeur externe", "—", "—", "20-40j × 500-800€"),
    ("Apple Developer", "0€", "99€/an", "99€/an"),
    ("Google Play", "25$ unique", "25$ unique", "25$ unique"),
    ("Firebase (notifications)", "Gratuit", "Gratuit", "Gratuit"),
    ("Tags NFC", "—", "~30€ (100 tags)", "—"),
    ("", "", "", ""),
    ("Total année 1", "~25$", "~180€", "10 000 - 32 000€"),
]
for ri, (label, *vals) in enumerate(rows_data):
    y = Inches(1.8) + Inches(ri * 0.5)
    bg = SURFACE if ri % 2 == 0 else SURFACE2
    is_total = "Total" in label
    rect(s, Inches(0.8), y, Inches(11.7), Inches(0.45), bg, ACCENT if is_total else None if is_total else BORDER)
    tf = textbox(s, Inches(0.9), y + Inches(0.05), Inches(3.2), Inches(0.35))
    p_add(tf, label, 10, is_total, TEXT)
    for vi, v in enumerate(vals):
        tf = textbox(s, cols_x[vi+1] + Inches(0.1), y + Inches(0.05), (Inches(2.4) if vi < 2 else Inches(2.4)) - Inches(0.2), Inches(0.35))
        c = ACCENT if "0€" in v or "Gratuit" in v or "~25" in v or "~180" in v else (RGBColor(0xDC, 0x26, 0x26) if "32" in v else TEXT_SEC)
        bold = any(x in v for x in ["~25", "~180", "10 ", "32"])
        p_add(tf, v, 10, bold, c, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════
s = new_slide()
set_bg(s, ACCENT_DARK)
circle(s, Inches(-0.5), Inches(-0.5), Inches(2.5), RGBColor(0x00, 0xAA, 0x88))
circle(s, Inches(11), Inches(5), Inches(3.5), RGBColor(0x00, 0x99, 0x77))

tf = textbox(s, Inches(1.5), Inches(1.5), Inches(10), Inches(1))
p_add(tf, "🎯 Ma recommandation", 40, True, WHITE, PP_ALIGN.LEFT, "Calibri Light")

tf = textbox(s, Inches(1.5), Inches(2.8), Inches(10), Inches(3))
p_add(tf, "Phase 1 — PWA (1 jour)", 24, True, ACCENT, name="Calibri Light")
p_add(tf, "Lance la démo fondation avec l'installation sur téléphone.", 14, False, RGBColor(0xBB, 0xF7, 0xE0))
p_add(tf, "", 10)
p_add(tf, "Phase 2 — Capacitor (3-5 jours) ← Recommandé", 24, True, AMBER, name="Calibri Light")
p_add(tf, "Le NFC natif est un argument différenciant fort pour La France s'engage.", 14, False, RGBColor(0xBB, 0xF7, 0xE0))
p_add(tf, "Tu gardes 95% du code Next.js, tu ajoutes juste le wrapper natif.", 14, False, RGBColor(0xBB, 0xF7, 0xE0))
p_add(tf, "", 10)
p_add(tf, "Phase 3 — React Native (4-8 semaines)", 24, True, RGBColor(0xCC, 0xCC, 0xCC), name="Calibri Light")
p_add(tf, "Uniquement si tu passes à 10 000+ utilisateurs.", 14, False, RGBColor(0xAA, 0xAA, 0xAA))

tf = textbox(s, Inches(1.5), Inches(6.2), Inches(10), Inches(0.6))
p_add(tf, "Prêt à démarrer la Phase 1 ?", 18, True, WHITE, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════
# SAVE & UPLOAD
# ═══════════════════════════════════════════════════════════════════════
output_path = "/root/projects/timebank-poc/Cours-App-Mobile-TimeHeroes.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
print(f"📊 Slides: {len(prs.slides)}")
