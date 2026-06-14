#!/usr/bin/env python3
"""Generate a light-themed TimeHeroes MBA pitch deck in .pptx format."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime, timezone

# ── Light color palette (from live site) ──
BG = RGBColor(0xF8, 0xF5, 0xF0)       # warm cream background
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)   # white cards
SURFACE2 = RGBColor(0xF0, 0xED, 0xE8) # light gray card bg
SURFACE3 = RGBColor(0xE8, 0xE5, 0xE0) # hover bg
ACCENT = RGBColor(0x00, 0xD4, 0xAA)    # teal
ACCENT_DARK = RGBColor(0x00, 0xA8, 0x89) # hover teal
TEXT = RGBColor(0x10, 0x10, 0x10)      # near-black
TEXT_SEC = RGBColor(0x5F, 0x63, 0x68)  # medium gray
TEXT_MUTED = RGBColor(0x9A, 0x9E, 0xA2) # light gray
BORDER = RGBColor(0xE0, 0xDD, 0xD8)    # light border
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF0, 0xA5, 0x00)
PURPLE = RGBColor(0x7C, 0x5C, 0xFC)
GREEN_DARK = RGBColor(0x0D, 0x54, 0x2B) # dark green accent

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ── Helpers ──

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines,
                          default_size=14, default_color=TEXT_SEC,
                          default_align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold = line_data, default_size, default_color, False
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = default_align
        p.space_after = Pt(4)
    return txBox

def add_tag(slide, left, top, text, fill=None):
    if fill is None:
        fill = SURFACE2
    w = Inches(0.4 + len(text) * 0.12)
    shape = add_rounded_rect(slide, left, top, max(w, Inches(1.8)), Inches(0.32), fill, BORDER)
    shape.adjustments[0] = 0.5
    add_textbox(slide, left + Inches(0.08), top + Inches(0.02), max(w, Inches(1.6)), Inches(0.28),
                text, font_size=10, color=ACCENT, alignment=PP_ALIGN.CENTER)
    return shape

# ── Layout constants ──
start_x = Inches(0.6)
gap = Inches(0.3)

# ══════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_textbox(slide, Inches(4.5), Inches(1.0), Inches(4.5), Inches(0.5),
            "⚡ STRATEGIC PROJECT", font_size=16, color=ACCENT,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.5), Inches(1.6), Inches(10.5), Inches(2.0),
            "TIMEHEROES", font_size=80, color=TEXT,
            alignment=PP_ALIGN.CENTER, font_name='Arial Black')

add_textbox(slide, Inches(1.5), Inches(3.4), Inches(10.5), Inches(1.0),
            "Une banque du temps moderne pour les héros du quotidien",
            font_size=22, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

tags_data = [("🎯 POC fonctionnel", 3.2), ("🏗️ 12 lots livrés", 5.7), ("🚀 timeheroes.fr", 7.8)]
for txt, x in tags_data:
    w = Inches(0.3 + len(txt) * 0.16)
    shape = add_rounded_rect(slide, Inches(x), Inches(4.6), w, Inches(0.38), WHITE, BORDER)
    shape.adjustments[0] = 0.5
    add_textbox(slide, Inches(x) + Inches(0.08), Inches(4.63), w, Inches(0.35),
                txt, font_size=12, color=ACCENT, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(3.5), Inches(6.0), Inches(6.5), Inches(0.5),
            "Ronald Mounien — MBA ESSEC Executive 2026", font_size=14, color=TEXT_MUTED,
            alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.5), "Pourquoi ce projet ?")
add_textbox(slide, Inches(0.6), Inches(1.0), Inches(9.0), Inches(1.2),
            "Le temps, dernière ressource vraiment égale",
            font_size=44, color=TEXT, font_name='Arial Black')

problems = [
    ("🤝", "Lien social en déclin",
     "Nos quartiers se fragmentent. 1 Français sur 3\nne connaît pas ses voisins. Les plateformes\nmarchandes ne font que consumer ce lien."),
    ("💸", "Barrière économique",
     "Aide informatique, soutien scolaire, bricolage —\nces services existent mais sont inaccessibles\nà ceux qui n'ont pas les moyens de payer."),
    ("📉", "Engagement non valorisé",
     "Des millions d'heures de bénévolat ne laissent\naucune trace. Pas de reconnaissance, pas de CV,\npas de valorisation professionnelle."),
]

card_w = Inches(3.8)
card_h = Inches(3.8)
card_y = Inches(2.6)

for i, (emoji, title, desc) in enumerate(problems):
    x = start_x + i * (card_w + gap)
    card = add_rounded_rect(slide, x, card_y, card_w, card_h, WHITE, BORDER)
    card.adjustments[0] = 0.04
    add_textbox(slide, x + Inches(0.3), card_y + Inches(0.3), Inches(0.6), Inches(0.6),
                emoji, font_size=32)
    add_textbox(slide, x + Inches(0.3), card_y + Inches(1.1), card_w - Inches(0.6), Inches(0.5),
                title, font_size=18, color=TEXT, bold=True)
    add_textbox(slide, x + Inches(0.3), card_y + Inches(1.7), card_w - Inches(0.6), Inches(1.8),
                desc, font_size=13, color=TEXT_SEC)

# ══════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.5), "La solution")
add_textbox(slide, Inches(0.6), Inches(1.0), Inches(9.0), Inches(1.0),
            "Une banque du temps — 1 heure = 1 TIME",
            font_size=44, color=TEXT, font_name='Arial Black')

add_textbox(slide, Inches(0.6), Inches(2.0), Inches(11.0), Inches(1.0),
            "TimeHeroes est une plateforme où chaque compétence devient une monnaie d'échange.\nPas d'argent. Du temps contre du temps. Avec des mécanismes modernes pour\nsécuriser, valoriser et gamifier l'entraide.",
            font_size=16, color=TEXT_SEC)

props = [
    ("🔒", "Escrow TIME",
     "Les TIME sont bloqués pendant la\nmission. Libérés uniquement après\nvalidation des deux parties."),
    ("✅", "QR Code & NFC",
     "Validation de complétion par QR\ncode scanné ou tag NFC. Zéro\nfriction entre héros."),
    ("🏆", "Gamification Hero",
     "XP, badges, niveaux, quêtes et\nrécompenses. Chaque heure\nd'entraide compte."),
]

for i, (emoji, title, desc) in enumerate(props):
    x = start_x + i * (card_w + gap)
    y = Inches(3.2)
    card = add_rounded_rect(slide, x, y, card_w, Inches(3.2), WHITE, BORDER)
    card.adjustments[0] = 0.04
    add_textbox(slide, x + Inches(card_w/2 - 0.3), y + Inches(0.3), Inches(0.7), Inches(0.6),
                emoji, font_size=36, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.6), Inches(0.5),
                title, font_size=18, color=TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.8), card_w - Inches(0.6), Inches(1.5),
                desc, font_size=14, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 4 — FEATURES
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.3), "Fonctionnalités")
add_textbox(slide, Inches(0.6), Inches(0.65), Inches(8.0), Inches(0.8),
            "Tout ce qu'un héros mérite d'avoir",
            font_size=36, color=TEXT, font_name='Arial Black')

features = [
    ("🏪", "Marketplace", "Recherche, filtres par catégories,\ndistance et localisation."),
    ("💰", "Wallet & TIME", "Portefeuille TIME avec historique,\ntransferts, mint initial."),
    ("📅", "Booking & Escrow", "Réservation, blocage en escrow,\nannulation, complétion."),
    ("⭐", "Réputation", "Système de notes et avis\npost-mission."),
    ("🆘", "Urgent Help", "Mode demande urgente. La\ncommunauté répond vite."),
    ("📍", "Local Heroes", "Géolocalisation déclarative.\nHéros dans ton quartier."),
    ("📱", "QR / NFC", "Validation physique de\ncomplétion par scan/tag."),
    ("🎮", "Gamification", "XP, badges, niveaux, quêtes\net récompenses."),
    ("🎯", "Pot Commun", "Financement solidaire des\nmissions communautaires."),
]

fw = Inches(3.85)
fh = Inches(1.35)
cols = 3
pad_x = Inches(0.35)
pad_y = Inches(0.25)
start_fx = Inches(0.6)
start_fy = Inches(1.6)

for i, (emoji, title, desc) in enumerate(features):
    col = i % cols
    row = i // cols
    x = start_fx + col * (fw + pad_x)
    y = start_fy + row * (fh + pad_y)
    card = add_rounded_rect(slide, x, y, fw, fh, WHITE, BORDER)
    card.adjustments[0] = 0.04
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.4), Inches(0.5),
                emoji, font_size=24)
    add_textbox(slide, x + Inches(0.7), y + Inches(0.15), Inches(2.8), Inches(0.4),
                title, font_size=15, color=TEXT, bold=True)
    add_textbox(slide, x + Inches(0.7), y + Inches(0.55), Inches(2.9), Inches(0.7),
                desc, font_size=12, color=TEXT_SEC)

# ══════════════════════════════════════════════
# SLIDE 5 — GAMIFICATION
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.5), "⭐ Différentiateur clé")
add_textbox(slide, Inches(0.6), Inches(1.0), Inches(9.0), Inches(1.0),
            "Gamification Hero : ton entraide devient ton XP",
            font_size=40, color=TEXT, font_name='Arial Black')

gitems = [
    ("⭐", "5 Niveaux", "Débutant → Héros Légendaire.\nChaque niveau débloque des privilèges."),
    ("🎖️", "Badges", "Badges thématiques pour les\nactions spéciales et l'engagement."),
    ("📋", "Quêtes", "Défis quotidiens et missions\nspéciales pour rester engagé."),
]

for i, (emoji, title, desc) in enumerate(gitems):
    x = start_x + i * (card_w + gap)
    y = Inches(2.4)
    card = add_rounded_rect(slide, x, y, card_w, Inches(2.8), WHITE, BORDER)
    card.adjustments[0] = 0.04
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.5), Inches(0.5),
                emoji, font_size=28)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.9), card_w - Inches(0.6), Inches(0.5),
                title, font_size=22, color=ACCENT, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.6), Inches(1.2),
                desc, font_size=14, color=TEXT_SEC)

# Highlight
hl = add_rounded_rect(slide, Inches(0.6), Inches(5.6), Inches(12.0), Inches(1.3), WHITE, BORDER)
hl.adjustments[0] = 0.04
add_multiline_textbox(slide, Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.2), [
    ("≠ Différence fondamentale :", 14, ACCENT, True),
    ("La plupart des time banks s'arrêtent au matching. TimeHeroes transforme chaque échange en expérience engageante —", 13, TEXT_SEC),
    ("comme un RPG social où plus tu aides, plus tu montes en grade. Et ça devient un CV de l'engagement.", 13, TEXT_SEC),
])

# ══════════════════════════════════════════════
# SLIDE 6 — USER JOURNEY
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.5), "Parcours utilisateur")
add_textbox(slide, Inches(0.6), Inches(1.0), Inches(9.0), Inches(0.8),
            "Du besoin → au héros en 5 minutes",
            font_size=40, color=TEXT, font_name='Arial Black')

steps_left = [
    ("1", "Je crée mon compte Hero",
     "Inscription rapide avec email. Ou connexion démo\n1-clic pour tester immédiatement."),
    ("2", "Je propose ou je cherche",
     "Je publie ce que je sais faire, ou je trouve\nquelqu'un pour m'aider. Filtres par catégorie."),
    ("3", "Je réserve en quelques clics",
     "Les TIME sont bloqués en escrow. Le héros\nest prévenu. La mission est officielle."),
]
steps_right = [
    ("4", "La mission se réalise",
     "Cours de guitare, aide déménagement, soutien\nscolaire… le temps s'échange."),
    ("5", "Je valide par QR/NFC",
     "Un scan. Un tag. Les TIME sont libérés.\nLes deux parties notent l'expérience."),
    ("6", "Je gagne XP + badges 🏆",
     "Ma réputation grandit. Mon niveau grimpe.\nMon CV d'engagement s'enrichit."),
]

def draw_steps(slide, steps, offset_x, start_y=Inches(2.2)):
    for i, (num, title, desc) in enumerate(steps):
        y = start_y + i * (Inches(1.3) + Inches(0.15))
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, offset_x, y, Inches(0.5), Inches(0.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = ACCENT
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = ACCENT
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        add_textbox(slide, offset_x + Inches(0.7), y - Inches(0.02), Inches(5.0), Inches(0.35),
                    title, font_size=16, color=TEXT, bold=True)
        add_textbox(slide, offset_x + Inches(0.7), y + Inches(0.35), Inches(5.0), Inches(0.7),
                    desc, font_size=12, color=TEXT_SEC)

draw_steps(slide, steps_left, Inches(0.6))
draw_steps(slide, steps_right, Inches(6.8))

# ══════════════════════════════════════════════
# SLIDE 7 — METRICS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.5), "Chiffres & Réalisation")
add_textbox(slide, Inches(0.6), Inches(1.0), Inches(9.0), Inches(0.8),
            "Du concept à la réalité en 2 mois",
            font_size=40, color=TEXT, font_name='Arial Black')

stats = [
    ("12", "Lots fonctionnels\nlivrés"),
    ("22+", "Routes /\npages"),
    ("10", "Catégories de\nservices"),
    ("5", "Niveaux\nHero"),
]

for i, (num, label) in enumerate(stats):
    x = start_x + i * (Inches(2.9) + Inches(0.25))
    y = Inches(2.2)
    card = add_rounded_rect(slide, x, y, Inches(2.9), Inches(2.0), WHITE, BORDER)
    card.adjustments[0] = 0.04
    add_textbox(slide, x, y + Inches(0.2), Inches(2.9), Inches(0.8),
                num, font_size=48, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER, font_name='Arial Black')
    add_textbox(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.5), Inches(0.8),
                label, font_size=13, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

details = [
    ("Lots livrés", "Auth & Wallet · Marketplace · Booking & Escrow ·\nRating & Reputation · Local Heroes · Urgent Help ·\nNFC Proof · Gamification · Demo Seed · Landing Page"),
    ("Stack", "Next.js 16 · TypeScript · Prisma 6 · SQLite ·\nTailwind v4 · NextAuth · Zod"),
]

for i, (title, content) in enumerate(details):
    x = start_x + i * (Inches(6.0) + Inches(0.3))
    y = Inches(4.6)
    card = add_rounded_rect(slide, x, y, Inches(6.0), Inches(2.2), WHITE, BORDER)
    card.adjustments[0] = 0.04
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), Inches(5.5), Inches(0.4),
                title, font_size=12, color=TEXT_MUTED)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.65), Inches(5.5), Inches(1.3),
                content, font_size=13, color=TEXT_SEC)

# ══════════════════════════════════════════════
# SLIDE 8 — ARCHITECTURE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.5), "Architecture & Déploiement")
add_textbox(slide, Inches(0.6), Inches(1.0), Inches(9.0), Inches(0.8),
            "Production ready",
            font_size=40, color=TEXT, font_name='Arial Black')
add_textbox(slide, Inches(0.6), Inches(1.7), Inches(9.0), Inches(0.5),
            "Une stack moderne, un déploiement VPS réel",
            font_size=16, color=TEXT_SEC)

stack = ["Next.js 16", "TypeScript", "Prisma 6", "SQLite", "Tailwind v4", "NextAuth", "Zod"]
for i, item in enumerate(stack):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + col * Inches(2.8)
    y = Inches(2.4) + row * Inches(0.42)
    shape = add_rounded_rect(slide, x, y, Inches(2.5), Inches(0.35), WHITE, BORDER)
    shape.adjustments[0] = 0.5
    add_textbox(slide, x + Inches(0.08), y + Inches(0.02), Inches(2.3), Inches(0.3),
                item, font_size=12, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)

flow_box = add_rounded_rect(slide, Inches(0.6), Inches(3.5), Inches(5.5), Inches(2.5), WHITE, BORDER)
flow_box.adjustments[0] = 0.04
add_textbox(slide, Inches(0.9), Inches(3.6), Inches(5.0), Inches(0.4),
            "Déploiement", font_size=14, color=ACCENT, bold=True)
deploy_lines = [
    "Client → HTTPS timeheroes.fr",
    "  → Caddy (SSL) → localhost:3096",
    "    → Next.js (programmatic server)",
    "      → Prisma → SQLite",
]
add_multiline_textbox(slide, Inches(0.9), Inches(4.1), Inches(5.0), Inches(1.6),
                      [(l, 13, TEXT_SEC) for l in deploy_lines])

demo_box = add_rounded_rect(slide, Inches(6.5), Inches(3.5), Inches(6.0), Inches(2.5), WHITE, ACCENT)
demo_box.adjustments[0] = 0.04
demo_box.line.width = Pt(2)
add_textbox(slide, Inches(6.8), Inches(3.7), Inches(5.5), Inches(0.4),
            "🌐  timeheroes.fr", font_size=20, color=TEXT, bold=True)
add_textbox(slide, Inches(6.8), Inches(4.2), Inches(5.5), Inches(0.4),
            "Hébergé sur VPS Hetzner · Domaine OVH · SSL Let's Encrypt",
            font_size=12, color=TEXT_SEC)

add_textbox(slide, Inches(6.8), Inches(4.8), Inches(5.5), Inches(0.3),
            "⚠️  Testez la démo", font_size=13, color=TEXT_MUTED)
add_textbox(slide, Inches(6.8), Inches(5.1), Inches(5.5), Inches(0.3),
            "Email : demo@timeheroes.fr", font_size=14, color=TEXT, bold=True)
add_textbox(slide, Inches(6.8), Inches(5.4), Inches(5.5), Inches(0.3),
            "Mot de passe : TimeHeroes2026!", font_size=14, color=TEXT, bold=True)
add_textbox(slide, Inches(6.8), Inches(5.8), Inches(5.5), Inches(0.3),
            "Ou utilisez le bouton ⚡ Connexion démo 1-clic",
            font_size=12, color=TEXT_SEC)

# ══════════════════════════════════════════════
# SLIDE 9 — ROADMAP
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_tag(slide, Inches(0.6), Inches(0.5), "Feuille de route")
add_textbox(slide, Inches(0.6), Inches(1.0), Inches(9.0), Inches(0.8),
            "La suite de l'aventure TimeHeroes",
            font_size=40, color=TEXT, font_name='Arial Black')

phases = [
    ("✅ Fait — POC (Juin 2026)", ACCENT,
     ["Auth, Wallet, Marketplace",
      "Booking & Escrow",
      "QR/NFC, Gamification",
      "Urgent Help, Local Heroes",
      "Landing page, seed data"]),
    ("🎯 En cours — MVP", ORANGE,
     ["Dashboard admin analytics",
      "Notifications temps réel",
      "Messagerie intégrée",
      "Onboarding enrichi",
      "PostgreSQL scale"]),
    ("🌟 Vision — Long terme", PURPLE,
     ["Records of Achievement",
      "Portfolio bénévole / CV",
      "Réseau coopératives ESS",
      "Application mobile",
      "API publique"]),
]

phase_w = Inches(3.8)
for i, (label, accent, items) in enumerate(phases):
    x = start_x + i * (phase_w + gap)
    y = Inches(2.2)
    card = add_rounded_rect(slide, x, y, phase_w, Inches(4.5), WHITE, BORDER)
    card.adjustments[0] = 0.04
    bar = add_rect(slide, x + Inches(0.05), y + Inches(0.05), phase_w - Inches(0.1), Inches(0.04), accent)
    add_textbox(slide, x + Inches(0.25), y + Inches(0.3), phase_w - Inches(0.5), Inches(0.35),
                label, font_size=13, color=accent, bold=True)
    for j, item in enumerate(items):
        iy = y + Inches(0.85) + j * Inches(0.55)
        add_textbox(slide, x + Inches(0.25), iy, phase_w - Inches(0.5), Inches(0.35),
                    f"→  {item}", font_size=13, color=TEXT_SEC)

# ══════════════════════════════════════════════
# SLIDE 10 — CTA / VISION
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG)

add_textbox(slide, Inches(2.0), Inches(1.2), Inches(9.5), Inches(0.6),
            "🌍 L'AVENIR DE L'ENTRAIDE", font_size=18, color=ACCENT,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2),
            "L'engagement citoyen devient\nun actif professionnel",
            font_size=48, color=TEXT, alignment=PP_ALIGN.CENTER, font_name='Arial Black')
add_multiline_textbox(slide, Inches(1.5), Inches(3.5), Inches(10.5), Inches(1.2), [
    ("TimeHeroes, c'est plus qu'une app. C'est un nouveau standard pour", 17, TEXT_SEC),
    ("reconnaître, valoriser et monétiser l'engagement de chacun.", 17, TEXT_SEC),
    ("Le temps est la seule vraie richesse. Échangeons-la.", 15, TEXT_MUTED),
], default_align=PP_ALIGN.CENTER)

cta = add_rounded_rect(slide, Inches(3.0), Inches(5.0), Inches(7.5), Inches(1.5), WHITE, ACCENT)
cta.adjustments[0] = 0.04
cta.line.width = Pt(2)
add_textbox(slide, Inches(3.3), Inches(5.1), Inches(7.0), Inches(0.5),
            "🚀  timeheroes.fr", font_size=28, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(3.3), Inches(5.6), Inches(7.0), Inches(0.4),
            "📧  demo@timeheroes.fr  •  Mot de passe : TimeHeroes2026!",
            font_size=14, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(6.6), Inches(10.5), Inches(0.6),
            "\"Le temps est la seule monnaie qui ne s'emprunte pas, ne s'imprime pas, ne se crée pas.\"",
            font_size=12, color=TEXT_MUTED, alignment=PP_ALIGN.CENTER)

# ── Save locally ──
output_path = "/root/projects/timebank-poc/TimeHeroes - MBA Pitch Deck.pptx"
prs.save(output_path)
print(f"✅ Saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
