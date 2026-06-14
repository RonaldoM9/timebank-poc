#!/usr/bin/env python3
"""Generate TimeHeroes pitch deck PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──
BG_DARK = RGBColor(0x0A, 0x0A, 0x0A)
BG_SURFACE = RGBColor(0x11, 0x11, 0x11)
BG_ELEVATED = RGBColor(0x18, 0x18, 0x18)
BG_CARD = RGBColor(0x15, 0x15, 0x15)
ACCENT = RGBColor(0x00, 0xD4, 0xAA)
ACCENT_DARK = RGBColor(0x00, 0xB8, 0x94)
TEXT_PRIMARY = RGBColor(0xF5, 0xF5, 0xF5)
TEXT_SECONDARY = RGBColor(0xA3, 0xA3, 0xA3)
TEXT_MUTED = RGBColor(0x73, 0x73, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x02, 0x06, 0x17)
CARD_BORDER = RGBColor(0x26, 0x26, 0x26)

# ── Helpers ──
def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14, color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, font_name='Calibri', line_spacing=1.5):
    """lines: list of (text, bold, color_override, font_size_override)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        text = line_info[0]
        bold = line_info[1] if len(line_info) > 1 else False
        clr = line_info[2] if len(line_info) > 2 else color
        fsize = line_info[3] if len(line_info) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fsize)
        p.font.bold = bold
        p.font.color.rgb = clr
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(fsize * 0.3)
    return txBox

def add_card(slide, left, top, width, height, accent_color=None):
    """Add a card with optional left accent bar."""
    card = add_rounded_rect(slide, left, top, width, height, BG_CARD)
    card.shadow.inherit = False
    if accent_color:
        add_shape(slide, left + Inches(0.06), top + Inches(0.15), Inches(0.05), height - Inches(0.3), accent_color)
    return card

def add_icon_circle(slide, left, top, size, icon_text, bg_color=ACCENT, text_color=BLACK):
    """Add a circle with emoji icon."""
    c = add_circle(slide, left, top, size, bg_color)
    add_textbox(slide, left, top, size, size, icon_text, font_size=int(size / Pt(1) * 0.5), 
                align=PP_ALIGN.CENTER, color=text_color)
    return c

def add_halftone_dots(slide):
    """Add decorative small dots pattern."""
    for x in range(0, 10, 2):
        for y in range(0, 6, 2):
            dot = add_circle(slide, Inches(x * 0.08), Inches(y * 0.08), 
                           Inches(0.03), RGBColor(0x00, 0xD4, 0xAA))

# ── Create Presentation ──
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, BG_DARK)

# Decorative circles
add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(2.5), RGBColor(0x00, 0xD4, 0xAA))
add_circle(slide, Inches(10), Inches(-1), Inches(3), RGBColor(0x00, 0xD4, 0xAA))
add_circle(slide, Inches(0.5), Inches(6), Inches(1.5), RGBColor(0x00, 0xD4, 0xAA))

# Subtitle label
add_textbox(slide, Inches(1), Inches(1.8), Inches(5), Inches(0.5), 
            "TIMEHEROES", font_size=11, bold=True, color=ACCENT, font_name='Arial Black')

# Main title
add_textbox(slide, Inches(1), Inches(2.3), Inches(11), Inches(1.2),
            "Échangez du temps,\npas de l'argent", font_size=56, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')

# Divider line
add_shape(slide, Inches(1), Inches(4.0), Inches(1.5), Inches(0.04), ACCENT)

# Subtitle
add_textbox(slide, Inches(1), Inches(4.3), Inches(10), Inches(0.8),
            "La plateforme de TimeBanking où chacun devient un super-héros du quotidien.\nDonnez du temps, recevez du temps, créez du lien.",
            font_size=18, color=TEXT_SECONDARY)

# Tagline
add_textbox(slide, Inches(1), Inches(5.2), Inches(8), Inches(0.5),
            "Nous sommes tous des super-héros du quotidien.",
            font_size=14, bold=True, color=ACCENT)

# Big Hero T on the right
add_textbox(slide, Inches(9.5), Inches(1.5), Inches(3), Inches(4),
            "T", font_size=180, bold=True, color=ACCENT, 
            align=PP_ALIGN.CENTER, font_name='Arial Black')

# Bottom bar
add_shape(slide, Inches(0), Inches(7.0), W, Inches(0.5), RGBColor(0x00, 0xD4, 0xAA))
add_textbox(slide, Inches(0.5), Inches(7.05), Inches(4), Inches(0.4),
            "POC — MBA ESSEC 2026 • Ronald Mounien", font_size=10, color=TEXT_MUTED)
add_textbox(slide, Inches(8), Inches(7.05), Inches(5), Inches(0.4),
            "timeheroes.fr", font_size=10, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Header
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.3),
            "LE PROBLÈME", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.8),
            "Pourquoi le timebanking est une nécessité", font_size=40, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.7), Inches(1.2), Inches(0.04), ACCENT)

# 3 problem cards
problems = [
    ("💰", "L'argent crée des barrières", "Beaucoup de services utiles restent inaccessibles parce qu'ils ont un coût. Pourtant, les compétences existent autour de nous."),
    ("🤝", "Le lien social s'affaiblit", "Dans les quartiers et les communautés, on se connaît moins, on s'entraide moins. Le potentiel de solidarité est sous-exploité."),
    ("📋", "L'engagement bénévole invisible", "Des centaines d'heures d'entraide ne laissent aucune trace — pas de CV, pas de portfolio, pas de reconnaissance professionnelle."),
]

card_w = Inches(3.6)
card_h = Inches(3.5)
gap = Inches(0.4)
start_x = Inches(0.8)
start_y = Inches(2.3)

for i, (icon, title, desc) in enumerate(problems):
    x = start_x + i * (card_w + gap)
    y = start_y
    
    card = add_rounded_rect(slide, x, y, card_w, card_h, BG_CARD)
    
    # Icon circle
    add_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), RGBColor(0x00, 0xD4, 0xAA))
    add_textbox(slide, x + Inches(0.3), y + Inches(0.35), Inches(0.6), Inches(0.5),
                icon, font_size=22, align=PP_ALIGN.CENTER)
    
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(1.1), card_w - Inches(0.6), Inches(0.5),
                title, font_size=18, bold=True, color=TEXT_PRIMARY)
    
    # Description
    add_textbox(slide, x + Inches(0.3), y + Inches(1.7), card_w - Inches(0.6), Inches(1.6),
                desc, font_size=13, color=TEXT_SECONDARY)

# Quote
quote_box = add_rounded_rect(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8), RGBColor(0x00, 0xD4, 0xAA))
add_textbox(slide, Inches(1.2), Inches(6.3), Inches(11), Inches(0.6),
            '"Et si votre temps avait autant de valeur que votre argent ?"',
            font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 3 — Solution / Concept
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Decorative
add_circle(slide, Inches(10), Inches(-0.3), Inches(2), RGBColor(0x00, 0xD4, 0xAA))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "LA SOLUTION", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.8),
            "TimeHeroes — Échangez du temps", font_size=40, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.7), Inches(1.2), Inches(0.04), ACCENT)

# Left: Concept
concept_x = Inches(0.8)
concept_y = Inches(2.3)
card_w2 = Inches(5.0)

# Concept card
add_rounded_rect(slide, concept_x, concept_y, card_w2, Inches(4.5), BG_CARD)
add_textbox(slide, concept_x + Inches(0.4), concept_y + Inches(0.3), card_w2 - Inches(0.8), Inches(0.4),
            "🎯 CONCEPT", font_size=13, bold=True, color=ACCENT, font_name='Arial Black')
add_multiline_textbox(slide, concept_x + Inches(0.4), concept_y + Inches(0.8), card_w2 - Inches(0.8), Inches(3.5), [
    ("Vous avez une compétence ? Proposez un service.", True, TEXT_PRIMARY, 16),
    ("Vous avez besoin d'aide ? Réservez avec vos TIME.", True, TEXT_PRIMARY, 16),
    ("", False, TEXT_MUTED, 8),
    ("1 heure de service = 1 TIME", True, ACCENT, 20),
    ("", False, TEXT_MUTED, 8),
    ("Une unité d'échange interne, simple et équitable.", False, TEXT_SECONDARY, 14),
], align=PP_ALIGN.LEFT)

# Right: Values
val_x = Inches(6.5)
val_y = Inches(2.3)
vals_w = Inches(6.0)

add_rounded_rect(slide, val_x, val_y, vals_w, Inches(4.5), BG_CARD)
add_textbox(slide, val_x + Inches(0.4), val_y + Inches(0.3), vals_w - Inches(0.8), Inches(0.4),
            "💎 VALEURS", font_size=13, bold=True, color=ACCENT, font_name='Arial Black')

values = [
    ("→ Équité", "1h de compétence = 1 TIME, quel que soit le service"),
    ("→ Confiance", "Système d'avis, escrow, badges de fiabilité"),
    ("→ Impact local", "Communautés, rayon de proximité"),
    ("→ Reconnaissance", "XP, niveaux, badges, futur passeport d'engagement"),
]
val_text_y = Inches(0.9)
for (label, desc) in values:
    add_textbox(slide, val_x + Inches(0.4), val_y + val_text_y, vals_w - Inches(0.8), Inches(0.3),
                label, font_size=15, bold=True, color=TEXT_PRIMARY)
    add_textbox(slide, val_x + Inches(0.4), val_y + val_text_y + Inches(0.3), vals_w - Inches(0.8), Inches(0.3),
                desc, font_size=12, color=TEXT_SECONDARY)
    val_text_y += Inches(0.8)

# Tags bottom
tags_y = Inches(6.1)
tags = ["Pas de crypto", "TIME ≠ monnaie", "100% solidaire", "Open data futur"]
for i, tag in enumerate(tags):
    tx = Inches(0.8) + i * Inches(2.8)
    tag_box = add_rounded_rect(slide, tx, tags_y, Inches(2.5), Inches(0.45), RGBColor(0x00, 0xD4, 0xAA))
    add_textbox(slide, tx, tags_y + Inches(0.05), Inches(2.5), Inches(0.35),
                tag, font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 4 — Key Features
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "FONCTIONNALITÉS", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.7),
            "Tout ce qu'un Hero peut faire", font_size=38, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.65), Inches(1.2), Inches(0.04), ACCENT)

features = [
    ("🛒", "Marketplace", "Proposez et trouvez des services par catégorie : numérique, bricolage, soutien scolaire, cuisine, etc."),
    ("🔒", "Escrow TIME", "TIME bloqués pendant la mission, libérés à la validation. Sécurité et confiance pour les deux parties."),
    ("⭐", "Avis & Réputation", "Système d'évaluation après chaque mission. Construisez votre réputation de Hero de confiance."),
    ("🎮", "Gamification Hero", "XP, niveaux (Rookie → Community Legend), badges de compétences et de confiance, quêtes."),
    ("👛", "Wallet TIME", "Solde, historique, transactions. Don de TIME au pot solidaire."),
    ("🌍", "Communautés locales", "Créez ou rejoignez des communautés géographiques ou thématiques."),
]

f_w = Inches(3.8)
f_h = Inches(2.0)
f_gap_x = Inches(0.3)
f_gap_y = Inches(0.25)
f_start_x = Inches(0.6)
f_start_y = Inches(2.1)

for i, (icon, title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = f_start_x + col * (f_w + f_gap_x)
    y = f_start_y + row * (f_h + f_gap_y)
    
    card = add_rounded_rect(slide, x, y, f_w, f_h, BG_CARD)
    
    # Icon
    add_circle(slide, x + Inches(0.25), y + Inches(0.25), Inches(0.45), RGBColor(0x00, 0xD4, 0xAA))
    add_textbox(slide, x + Inches(0.25), y + Inches(0.28), Inches(0.45), Inches(0.4),
                icon, font_size=16, align=PP_ALIGN.CENTER)
    
    add_textbox(slide, x + Inches(0.85), y + Inches(0.25), f_w - Inches(1.1), Inches(0.35),
                title, font_size=16, bold=True, color=TEXT_PRIMARY)
    add_textbox(slide, x + Inches(0.85), y + Inches(0.65), f_w - Inches(1.1), Inches(1.2),
                desc, font_size=11, color=TEXT_SECONDARY)

# ════════════════════════════════════════════════════════════
# SLIDE 5 — Levels & XP System
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_circle(slide, Inches(-0.5), Inches(1), Inches(2), RGBColor(0x00, 0xD4, 0xAA))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "GAMIFICATION", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.7),
            "Système Hero — XP, Niveaux & Progression", font_size=38, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.65), Inches(1.2), Inches(0.04), ACCENT)

# Level progression table
levels = [
    ("Rookie Hero 🆕", "0 XP", "Nouvel utilisateur"),
    ("Active Hero ⚡", "100 XP", "Premières actions"),
    ("Local Hero 🌍", "300 XP", "Utilisateur engagé"),
    ("Guardian Hero 🛡️", "700 XP", "Fiable et régulier"),
    ("Community Legend 👑", "1500 XP", "Contributeur majeur"),
]

table_x = Inches(0.8)
table_y = Inches(2.2)
col_w = [Inches(3.2), Inches(1.5), Inches(3.5)]

# Table header
add_rounded_rect(slide, table_x, table_y, sum([c.inches for c in col_w]) * 914400, Inches(0.45), RGBColor(0x00, 0xD4, 0xAA))
hdr_labels = ["NIVEAU", "XP", "SIGNIFICATION"]
hdr_x = table_x
for i, (label, w) in enumerate(zip(hdr_labels, col_w)):
    add_textbox(slide, hdr_x + Inches(0.15), table_y + Inches(0.05), w - Inches(0.3), Inches(0.35),
                label, font_size=10, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    hdr_x += w

# Table rows
for i, (name, xp, meaning) in enumerate(levels):
    ry = table_y + Inches(0.55) + i * Inches(0.5)
    bg = RGBColor(0x00, 0xD4, 0xAA) if i % 2 == 0 else RGBColor(0x00, 0x00, 0x00)
    
    rw = table_x
    add_textbox(slide, rw + Inches(0.15), ry + Inches(0.05), col_w[0] - Inches(0.3), Inches(0.35),
                name, font_size=16, bold=True, color=TEXT_PRIMARY)
    rw += col_w[0]
    add_textbox(slide, rw + Inches(0.15), ry + Inches(0.05), col_w[1] - Inches(0.3), Inches(0.35),
                xp, font_size=16, bold=True, color=ACCENT)
    rw += col_w[1]
    add_textbox(slide, rw + Inches(0.15), ry + Inches(0.05), col_w[2] - Inches(0.3), Inches(0.35),
                meaning, font_size=13, color=TEXT_SECONDARY)

# XP Rules box
rules_x = Inches(7.5)
rules_y = Inches(2.2)
rules_w = Inches(5.3)
rules_h = Inches(3.0)

add_rounded_rect(slide, rules_x, rules_y, rules_w, rules_h, BG_CARD)
add_textbox(slide, rules_x + Inches(0.3), rules_y + Inches(0.2), rules_w - Inches(0.6), Inches(0.4),
            "⚡ XP BARÈME", font_size=13, bold=True, color=ACCENT, font_name='Arial Black')

xp_items = [
    "Mission terminée (aidant)    +50 XP",
    "Mission terminée (bénéf.)    +20 XP",
    "Avis positif reçu              +30 XP",
    "Don de TIME                    +30 XP",
    "Profil complété                +20 XP",
    "Première mission publiée      +20 XP",
]

xp_text_y = Inches(0.7)
for item in xp_items:
    add_textbox(slide, rules_x + Inches(0.3), rules_y + xp_text_y, rules_w - Inches(0.6), Inches(0.28),
                item, font_size=12, color=TEXT_SECONDARY)
    xp_text_y += Inches(0.32)

# Key rule box
rule_y = Inches(5.5)
add_rounded_rect(slide, Inches(0.8), rule_y, Inches(12), Inches(0.7), RGBColor(0x00, 0xD4, 0xAA))
add_multiline_textbox(slide, Inches(1.2), rule_y + Inches(0.1), Inches(11.5), Inches(0.5), [
    ("⚠️ RÈGLE PRODUIT : XP ≠ TIME", True, ACCENT, 16),
    (" L'XP mesure la progression. Le TIME reste la devise utilisée pour les échanges. L'XP ne se dépense pas, ne se transfère pas.", False, TEXT_SECONDARY, 13),
])

# Tags
add_textbox(slide, Inches(0.8), Inches(6.6), Inches(12), Inches(0.4),
            "🎖️ Badges disponibles : First Mission  •  Helping Hand  •  Local Hero  •  Time Giver  •  Generous Hero  •  Trusted Hero  •  Reliable Hero  •  +10 autres",
            font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 6 — User Journey
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "PARCOURS", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.7),
            "Le parcours Hero", font_size=38, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.65), Inches(1.2), Inches(0.04), ACCENT)

journey = [
    ("1️⃣", "Inscription → 10 TIME", "Créez votre compte, complétez votre profil, recevez vos premiers TIME."),
    ("2️⃣", "Publier ou réserver", "Devenez aidant en proposant votre compétence, ou bénéficiaire."),
    ("3️⃣", "Mission en escrow", "Les TIME sont bloqués pendant la mission. Sécurité garantie."),
    ("4️⃣", "Validation + avis", "Terminez la mission, validez, libérez les TIME, laissez un avis."),
    ("5️⃣", "Progression Hero", "Gagnez XP, débloquez des badges, montez en niveau."),
    ("6️⃣", "Donner du TIME", "Soutenez le pot solidaire — chaque don est reconnu."),
]

# Timeline visual - vertical connected steps
timeline_x = Inches(2.5)
step_y = Inches(2.2)
step_h = Inches(0.75)
step_gap = Inches(0.15)

# Vertical line
add_shape(slide, timeline_x + Inches(0.02), step_y, Inches(0.03), Inches(5.5), RGBColor(0x00, 0xD4, 0xAA))

for i, (num, title, desc) in enumerate(journey):
    sy = step_y + i * (step_h + step_gap)
    
    # Circle marker
    add_circle(slide, timeline_x - Inches(0.15), sy + Inches(0.05), Inches(0.35), ACCENT)
    add_textbox(slide, timeline_x - Inches(0.15), sy + Inches(0.08), Inches(0.35), Inches(0.3),
                num, font_size=12, align=PP_ALIGN.CENTER, color=BLACK)
    
    # Card content
    cx = timeline_x + Inches(0.4)
    add_rounded_rect(slide, cx, sy, Inches(8), step_h, BG_CARD)
    
    # Title
    add_textbox(slide, cx + Inches(0.2), sy + Inches(0.05), Inches(5), Inches(0.3),
                title, font_size=14, bold=True, color=TEXT_PRIMARY)
    # Description
    add_textbox(slide, cx + Inches(0.2), sy + Inches(0.35), Inches(7), Inches(0.3),
                desc, font_size=11, color=TEXT_SECONDARY)

# ════════════════════════════════════════════════════════════
# SLIDE 7 — Key Metrics / Stats
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_circle(slide, Inches(11), Inches(-0.5), Inches(2), RGBColor(0x00, 0xD4, 0xAA))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "CHIFFRES", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.7),
            "TimeHeroes en chiffres", font_size=38, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.65), Inches(1.2), Inches(0.04), ACCENT)

# Big stat cards
stats = [
    ("10", "LOTS", "Du MVP à la gamification"),
    ("17", "BADGES", "Compétences, confiance & engagement"),
    ("5", "NIVEAUX", "Rookie → Community Legend"),
    ("6", "CATÉGORIES", "Numérique, brico, cuisine..."),
    ("29", "ROUTES", "Pages fonctionnelles"),
]

stat_w = Inches(2.2)
stat_h = Inches(2.2)
stat_start_x = Inches(0.4)
stat_gap = Inches(0.2)
stat_y = Inches(2.2)

for i, (num, label, desc) in enumerate(stats):
    sx = stat_start_x + i * (stat_w + stat_gap)
    card = add_rounded_rect(slide, sx, stat_y, stat_w, stat_h, BG_CARD)
    
    add_textbox(slide, sx, stat_y + Inches(0.3), stat_w, Inches(0.8),
                num, font_size=48, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, font_name='Arial Black')
    add_textbox(slide, sx, stat_y + Inches(1.2), stat_w, Inches(0.3),
                label, font_size=10, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
    add_textbox(slide, sx, stat_y + Inches(1.5), stat_w, Inches(0.5),
                desc, font_size=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# Tech stack card
stack_x = Inches(0.8)
stack_y = Inches(4.7)
stack_w = Inches(5.8)
stack_h = Inches(2.3)

add_rounded_rect(slide, stack_x, stack_y, stack_w, stack_h, BG_CARD)
add_textbox(slide, stack_x + Inches(0.3), stack_y + Inches(0.2), stack_w - Inches(0.6), Inches(0.3),
            "⚙️ TECH STACK", font_size=12, bold=True, color=ACCENT, font_name='Arial Black')

techs = [("Next.js 16", "Frontend + API Routes"), ("TypeScript", "Strict mode"), 
         ("Prisma 6", "ORM"), ("SQLite", "Database"),
         ("Tailwind v4", "CSS / Design system"), ("NextAuth", "Authentication")]

for i, (tech, role) in enumerate(techs):
    col = i % 2
    row = i // 2
    tx = stack_x + Inches(0.3) + col * Inches(2.7)
    ty = stack_y + Inches(0.6) + row * Inches(0.5)
    add_textbox(slide, tx, ty, Inches(2.5), Inches(0.22),
                f"▸ {tech}", font_size=12, bold=True, color=TEXT_PRIMARY)
    add_textbox(slide, tx + Inches(1.3), ty, Inches(1.2), Inches(0.22),
                role, font_size=10, color=TEXT_MUTED)

# Design card
design_x = Inches(7.0)
design_y = Inches(4.7)
design_w = Inches(5.8)
design_h = Inches(2.3)

add_rounded_rect(slide, design_x, design_y, design_w, design_h, BG_CARD)
add_textbox(slide, design_x + Inches(0.3), design_y + Inches(0.2), design_w - Inches(0.6), Inches(0.3),
            "🎨 DESIGN SYSTEM", font_size=12, bold=True, color=ACCENT, font_name='Arial Black')

design_items = ["Dark premium", "Accent #00d4aa", "Comics moderne", 
                "Halftone dots", "Burst shapes", "UI lisible"]

for i, item in enumerate(design_items):
    col = i % 3
    row = i // 3
    dx = design_x + Inches(0.3) + col * Inches(1.8)
    dy = design_y + Inches(0.6) + row * Inches(0.5)
    tag = add_rounded_rect(slide, dx, dy, Inches(1.6), Inches(0.35), RGBColor(0x00, 0xD4, 0xAA))
    add_textbox(slide, dx, dy + Inches(0.03), Inches(1.6), Inches(0.3),
                item, font_size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 8 — Architecture
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "ARCHITECTURE", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.7),
            "Architecture technique", font_size=38, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.65), Inches(1.2), Inches(0.04), ACCENT)

# Architecture flow diagram - using shapes
flow_w = Inches(1.8)
flow_h = Inches(0.7)
flow_gap = Inches(0.15)
flow_y = Inches(2.3)

flows = [
    ("🌐", "Navigateur"),
    ("⬇️", "HTTPS SSL"),
    ("🔄", "Caddy proxy"),
    ("⬢", "Next.js"),
    ("🗄️", "SQLite"),
]

for i, (icon, name) in enumerate(flows):
    fx = Inches(0.8) + i * (flow_w + flow_gap)
    card = add_rounded_rect(slide, fx, flow_y, flow_w, flow_h, BG_CARD)
    add_textbox(slide, fx, flow_y + Inches(0.1), flow_w, Inches(0.25),
                icon, font_size=18, align=PP_ALIGN.CENTER)
    add_textbox(slide, fx, flow_y + Inches(0.35), flow_w, Inches(0.25),
                name, font_size=11, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
    
    # Arrow between blocks
    if i < len(flows) - 1:
        ax = fx + flow_w
        add_textbox(slide, ax, flow_y + Inches(0.15), flow_gap, Inches(0.3),
                    "→", font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Deployment details card
deploy_x = Inches(0.8)
deploy_y = Inches(3.4)
deploy_w = Inches(11.5)
deploy_h = Inches(1.5)

add_rounded_rect(slide, deploy_x, deploy_y, deploy_w, deploy_h, BG_CARD)
add_textbox(slide, deploy_x + Inches(0.3), deploy_y + Inches(0.15), deploy_w - Inches(0.6), Inches(0.3),
            "🚀 DÉPLOIEMENT", font_size=13, bold=True, color=ACCENT, font_name='Arial Black')

deploy_info = [
    ("Hébergement :", "VPS Hetzner (204.168.193.43)"),
    ("Domaine :", "timeheroes.fr (OVH)"),
    ("SSL :", "Let's Encrypt via Caddy"),
    ("Service :", "Node.js via systemd (auto-reboot)"),
    ("Ports :", "80/443 (public), 3096 (interne)"),
]

for i, (label, val) in enumerate(deploy_info):
    col = i % 3
    row = i // 3
    dx = deploy_x + Inches(0.3) + col * Inches(3.8)
    dy = deploy_y + Inches(0.5) + row * Inches(0.45)
    add_textbox(slide, dx, dy, Inches(3.5), Inches(0.2),
                label, font_size=11, bold=True, color=TEXT_SECONDARY)
    add_textbox(slide, dx, dy + Inches(0.2), Inches(3.5), Inches(0.2),
                val, font_size=11, color=TEXT_PRIMARY)

# Test accounts
acct_y = Inches(5.3)
add_rounded_rect(slide, Inches(0.8), acct_y, Inches(11.5), Inches(0.8), RGBColor(0x00, 0xD4, 0xAA))
add_textbox(slide, Inches(1.2), acct_y + Inches(0.1), Inches(11), Inches(0.25),
            "🔑 COMPTES TEST", font_size=12, bold=True, color=ACCENT)
add_textbox(slide, Inches(1.2), acct_y + Inches(0.4), Inches(11), Inches(0.3),
            "ronald.test1@timebank.local / Lot2bTest2025!   •   alice.test1@timebank.local / Lot2bTest2025!",
            font_size=12, color=TEXT_SECONDARY)

# URL bar
url_y = Inches(6.3)
add_rounded_rect(slide, Inches(0.8), url_y, Inches(11.5), Inches(0.5), RGBColor(0x00, 0xD4, 0xAA))
add_textbox(slide, Inches(0.8), url_y + Inches(0.08), Inches(11.5), Inches(0.35),
            "🌐 https://timeheroes.fr",
            font_size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
# SLIDE 9 — Roadmap
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "ROADMAP", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.7),
            "Roadmap & Vision", font_size=38, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.65), Inches(1.2), Inches(0.04), ACCENT)

# Phase 1 — Done
ph1_y = Inches(2.2)
add_rounded_rect(slide, Inches(0.8), ph1_y, Inches(5.5), Inches(2.2), RGBColor(0x00, 0xD4, 0xAA))
add_textbox(slide, Inches(1.1), ph1_y + Inches(0.15), Inches(5), Inches(0.3),
            "PHASE 1 ✅ — POC LIVRÉ", font_size=14, bold=True, color=ACCENT, font_name='Arial Black')

ph1_items = [
    "✅ Lot 1-6 : Auth, services, wallet, marketplace, escrow",
    "✅ Lot 7 : Admin / NFC Proof of Completion",
    "✅ Lot 8 : Preuve de complétion NFC",
    "✅ Lot 9 : Notifications",
    "✅ Lot 10 : Gamification Hero (XP, badges, niveaux, quêtes)",
]

ph1_y2 = Inches(0.55)
for item in ph1_items:
    add_textbox(slide, Inches(1.1), ph1_y + ph1_y2, Inches(5), Inches(0.25),
                item, font_size=10, color=TEXT_PRIMARY)
    ph1_y2 += Inches(0.3)

# Phase 2 — Next
ph2_x = Inches(6.8)
add_rounded_rect(slide, ph2_x, ph1_y, Inches(5.5), Inches(2.2), BG_CARD)
add_textbox(slide, ph2_x + Inches(0.3), ph1_y + Inches(0.15), Inches(5), Inches(0.3),
            "PHASE 2 🎯 — PROCHAINES ÉTAPES", font_size=14, bold=True, color=ACCENT, font_name='Arial Black')

ph2_items = [
    "Badges saisonniers & challenges",
    "Carte d'impact partageable",
    "Admin badges complet",
    "Comptes de démo pour le pitch",
    "Tests utilisateurs réels",
]

ph2_y = Inches(0.55)
for item in ph2_items:
    add_textbox(slide, ph2_x + Inches(0.3), ph1_y + ph2_y, Inches(5), Inches(0.25),
                f"▸ {item}", font_size=11, color=TEXT_SECONDARY)
    ph2_y += Inches(0.3)

# Phase 3 — Vision
ph3_y = Inches(4.7)
add_rounded_rect(slide, Inches(0.8), ph3_y, Inches(11.5), Inches(2.0), RGBColor(0x00, 0xD4, 0xAA))
add_textbox(slide, Inches(1.1), ph3_y + Inches(0.15), Inches(11), Inches(0.3),
            "PHASE 3 🌟 — VISION LONG TERME", font_size=14, bold=True, color=ACCENT, font_name='Arial Black')

ph3_items = [
    ("📜", "Records of Achievement — Passeport d'engagement numérique"),
    ("🏛️", "ESS & Employabilité — L'engagement citoyen devient un actif professionnel"),
    ("💼", "Portfolio bénévole — CV valorisant les missions, badges et compétences acquises"),
    ("🌍", "Écosystème communautaire — Réseau de coopératives TimeBank ESS"),
]

ph3_y2 = Inches(0.55)
for icon, text in ph3_items:
    add_textbox(slide, Inches(1.1), ph3_y + ph3_y2, Inches(11), Inches(0.28),
                f"{icon}  {text}", font_size=12, color=TEXT_PRIMARY)
    ph3_y2 += Inches(0.34)

# ════════════════════════════════════════════════════════════
# SLIDE 10 — Vision / CTA
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_DARK)

# Decorative
add_circle(slide, Inches(-0.5), Inches(-0.5), Inches(2.5), RGBColor(0x00, 0xD4, 0xAA))
add_circle(slide, Inches(10.5), Inches(5.5), Inches(2), RGBColor(0x00, 0xD4, 0xAA))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.3),
            "VISION", font_size=10, bold=True, color=ACCENT, font_name='Arial Black')
add_textbox(slide, Inches(0.8), Inches(1.0), Inches(12), Inches(0.7),
            "Et après ? La grande vision", font_size=38, bold=True, color=TEXT_PRIMARY, font_name='Arial Black')
add_shape(slide, Inches(0.8), Inches(1.65), Inches(1.2), Inches(0.04), ACCENT)

# Two vision cards
v_w = Inches(5.3)
v_h = Inches(2.5)
v_x = Inches(0.8)
v_y = Inches(2.3)

# Card 1
add_rounded_rect(slide, v_x, v_y, v_w, v_h, BG_CARD)
add_textbox(slide, v_x + Inches(0.3), v_y + Inches(0.5), v_w - Inches(0.6), Inches(0.6),
            "📜", font_size=36, align=PP_ALIGN.CENTER)
add_textbox(slide, v_x + Inches(0.3), v_y + Inches(1.2), v_w - Inches(0.6), Inches(0.4),
            "Records of Achievement", font_size=20, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, v_x + Inches(0.3), v_y + Inches(1.6), v_w - Inches(0.6), Inches(0.7),
            "Un passeport d'engagement qui capture chaque mission, badge et heure donnée. Utilisable comme portfolio, CV bénévole et attestation d'impact.",
            font_size=12, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

# Card 2
v_x2 = Inches(6.7)
add_rounded_rect(slide, v_x2, v_y, v_w, v_h, BG_CARD)
add_textbox(slide, v_x2 + Inches(0.3), v_y + Inches(0.5), v_w - Inches(0.6), Inches(0.6),
            "🏛️", font_size=36, align=PP_ALIGN.CENTER)
add_textbox(slide, v_x2 + Inches(0.3), v_y + Inches(1.2), v_w - Inches(0.6), Inches(0.4),
            "ESS & Employabilité", font_size=20, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, v_x2 + Inches(0.3), v_y + Inches(1.6), v_w - Inches(0.6), Inches(0.7),
            "Intégration avec les structures de l'Économie Sociale et Solidaire, les entreprises et les collectivités. L'engagement citoyen devient un actif professionnel reconnu.",
            font_size=12, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

# CTA section
cta_y = Inches(5.3)
add_rounded_rect(slide, Inches(0.8), cta_y, Inches(11.5), Inches(1.5), RGBColor(0x00, 0xD4, 0xAA))

add_textbox(slide, Inches(0.8), cta_y + Inches(0.2), Inches(11.5), Inches(0.5),
            "TimeHeroes est prêt pour la démonstration",
            font_size=24, bold=True, color=TEXT_PRIMARY, align=PP_ALIGN.CENTER)
add_textbox(slide, Inches(0.8), cta_y + Inches(0.7), Inches(11.5), Inches(0.35),
            "POC fonctionnel — 10 lots livrés — En production sur timeheroes.fr",
            font_size=14, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

# CTA button visual
btn = add_rounded_rect(slide, Inches(4.5), cta_y + Inches(1.05), Inches(3), Inches(0.4), ACCENT)
add_textbox(slide, Inches(4.5), cta_y + Inches(1.08), Inches(3), Inches(0.35),
            "🚀  timeheroes.fr", font_size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Footer
add_textbox(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.35),
            "Ronald Mounien — MBA ESSEC Executive 2026 | Projet Strategic Project : TimeBanking",
            font_size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = "/root/projects/timebank-poc/TimeHeroes_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
print(f"Size: {os.path.getsize(output_path) / 1024:.1f} KB")
