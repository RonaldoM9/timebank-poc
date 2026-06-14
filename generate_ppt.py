#!/usr/bin/env python3
"""Generate TimeHeroes MBA Pitch Deck as PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Constants ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette (light theme)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xFA, 0xFA, 0xFA)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x1A)
TEXT_BODY = RGBColor(0x4A, 0x4A, 0x4A)
TEXT_MUTED = RGBColor(0x8A, 0x8A, 0x8A)
TEAL = RGBColor(0x00, 0xD4, 0xAA)
TEAL_DARK = RGBColor(0x00, 0xA0, 0x80)
TEAL_LIGHT = RGBColor(0xE0, 0xFA, 0xF5)
TEAL_GLOW = RGBColor(0xF0, 0xFF, 0xFB)
ACCENT_ORANGE = RGBColor(0xF0, 0xA5, 0x00)
ACCENT_PURPLE = RGBColor(0x7C, 0x5C, 0xFC)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]  # blank

def add_bg(slide, color=WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None, corner_radius=None):
    """Add a rectangle shape."""
    if corner_radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        if corner_radius:
            shape.adjustments[0] = corner_radius
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height):
    """Add a text box and return text frame for rich editing."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_paragraph(tf, text, font_size=14, bold=False, color=TEXT_BODY, alignment=PP_ALIGN.LEFT, font_name='Calibri', space_before=0, space_after=0):
    """Add a paragraph to a text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p

def add_run(p, text, font_size=None, bold=False, color=None, font_name=None):
    """Add a run to a paragraph."""
    run = p.add_run()
    run.text = text
    if font_size: run.font.size = Pt(font_size)
    if bold: run.font.bold = True
    if color: run.font.color.rgb = color
    if font_name: run.font.name = font_name
    return run

def add_teal_bar(slide, left=Inches(0), top=Inches(0), width=None, height=Pt(6)):
    """Add a teal accent bar."""
    w = width or SLIDE_W
    return add_shape(slide, left, top, w, height, fill_color=TEAL)

def add_slide_number(slide, num, total=10):
    """Add slide number at bottom right."""
    add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.3), Inches(0.4),
                f"{num}/{total}", font_size=10, color=TEXT_MUTED, alignment=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, OFF_WHITE)
add_teal_bar(slide, height=Pt(8))

# Title
add_textbox(slide, Inches(1.2), Inches(1.5), Inches(11), Inches(1.2),
            "TIMEHEROES", font_size=72, bold=True, color=TEXT_DARK, alignment=PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide, Inches(1.2), Inches(2.7), Inches(10), Inches(0.6),
            "Une banque du temps moderne pour les héros du quotidien",
            font_size=28, bold=False, color=TEXT_BODY, alignment=PP_ALIGN.LEFT)

# Teal accent line
add_shape(slide, Inches(1.2), Inches(3.5), Inches(2.5), Pt(4), fill_color=TEAL)

# Tags row
tf = add_rich_textbox(slide, Inches(1.2), Inches(3.9), Inches(10), Inches(0.5))
add_paragraph(tf, "🎯 POC fonctionnel    🏗️ 12 lots livrés    🚀 timeheroes.fr",
              font_size=14, color=TEAL_DARK, bold=True)

# Description
add_textbox(slide, Inches(1.2), Inches(4.6), Inches(8), Inches(1.2),
            "Une plateforme d'échange de services où chaque compétence devient une monnaie.\n"
            "Pas d'argent. Du temps contre du temps.",
            font_size=18, color=TEXT_BODY)

# Presenter
add_shape(slide, Inches(1.2), Inches(6.2), Inches(4), Pt(1), fill_color=BORDER_GRAY)
add_textbox(slide, Inches(1.2), Inches(6.4), Inches(6), Inches(0.5),
            "Ronald Mounien — MBA ESSEC Executive 2026",
            font_size=16, color=TEXT_MUTED)

# Right side decorative block
add_shape(slide, Inches(9.5), Inches(0.8), Inches(3.2), Inches(3.2), fill_color=TEAL_LIGHT, corner_radius=0.1)
add_textbox(slide, Inches(10.0), Inches(1.5), Inches(2.2), Inches(1.8),
            "PITCH\nDECK", font_size=42, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
            "POURQUOI CE PROJET ?", font_size=11, bold=True, color=TEAL_DARK, alignment=PP_ALIGN.LEFT)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "Le temps, dernière ressource vraiment égale",
            font_size=44, bold=True, color=TEXT_DARK)

# 3 problem cards
problems = [
    ("🤝", "Lien social en déclin", "1 Français sur 3 ne connaît pas ses voisins.\nLes plateformes marchandes ne font que\nconsommer ce lien."),
    ("💸", "Barrière économique", "Aide informatique, soutien scolaire, bricolage\n— ces services sont inaccessibles à ceux\nqui n'ont pas les moyens de payer."),
    ("📉", "Engagement non valorisé", "Des millions d'heures de bénévolat ne\nlaissent aucune trace. Pas de CV,\npas de valorisation professionnelle."),
]

card_w = Inches(3.6)
card_h = Inches(4.0)
card_gap = Inches(0.4)
start_x = Inches(0.8)

for i, (emoji, title, desc) in enumerate(problems):
    x = start_x + i * (card_w + card_gap)
    y = Inches(2.2)
    
    # Card background
    add_shape(slide, x, y, card_w, card_h, fill_color=LIGHT_GRAY, corner_radius=0.08)
    
    # Emoji
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.8), Inches(0.6),
                emoji, font_size=36, alignment=PP_ALIGN.LEFT)
    
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.6), Inches(0.5),
                title, font_size=20, bold=True, color=TEXT_DARK)
    
    # Description
    add_textbox(slide, x + Inches(0.3), y + Inches(1.6), card_w - Inches(0.6), Inches(2.0),
                desc, font_size=14, color=TEXT_BODY)

# ══════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
            "LA SOLUTION", font_size=11, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Une banque du temps  1 heure = 1 TIME",
            font_size=44, bold=True, color=TEXT_DARK)

add_textbox(slide, Inches(0.8), Inches(1.7), Inches(10), Inches(0.7),
            "TimeHeroes est une plateforme où chaque compétence devient une monnaie d'échange.\n"
            "Pas d'argent. Du temps contre du temps.",
            font_size=16, color=TEXT_BODY)

# 3 solution cards
solutions = [
    ("🔒", "Escrow TIME", "Les TIME sont bloqués pendant\nla mission. Libérés uniquement\naprès validation des deux parties."),
    ("✅", "QR Code & NFC", "Validation de complétion par QR\ncode scanné ou tag NFC.\nZéro friction entre héros."),
    ("🏆", "Gamification Hero", "XP, badges, niveaux, quêtes et\nrécompenses. Chaque heure\nd'entraide compte."),
]

for i, (emoji, title, desc) in enumerate(solutions):
    x = start_x + i * (card_w + card_gap)
    y = Inches(2.7)
    
    # Card with teal accent border
    add_shape(slide, x, y, card_w, Inches(3.8), fill_color=WHITE, line_color=TEAL, line_width=Pt(1.5), corner_radius=0.08)
    
    # Emoji
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.8), Inches(0.6),
                emoji, font_size=36)
    
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.6), Inches(0.5),
                title, font_size=20, bold=True, color=TEXT_DARK)
    
    # Description
    add_textbox(slide, x + Inches(0.3), y + Inches(1.6), card_w - Inches(0.6), Inches(1.8),
                desc, font_size=14, color=TEXT_BODY)

# ══════════════════════════════════════════════
# SLIDE 4 — KEY FEATURES
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, OFF_WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(4), Inches(0.4),
            "FONCTIONNALITÉS", font_size=11, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7),
            "Tout ce qu'un héros mérite d'avoir",
            font_size=40, bold=True, color=TEXT_DARK)

features = [
    ("🏪", "Marketplace", "Recherche, filtres par catégories,\ndistance et localisation."),
    ("💰", "Wallet & TIME", "Portefeuille TIME avec historique,\ntransferts, mint initial."),
    ("📅", "Booking & Escrow", "Réservation, blocage, annulation,\ncomplétion. Vie de la mission."),
    ("⭐", "Réputation", "Système de notes et avis\npost-mission."),
    ("🆘", "Urgent Help", "Mode demande urgente.\nLa communauté répond."),
    ("📍", "Local Heroes", "Géolocalisation. Trouve les héros\ndans ton quartier."),
    ("📱", "QR / NFC", "Validation physique de complétion.\nTechnologie NFC embarquée."),
    ("🎮", "Gamification", "XP, badges, niveaux, quêtes.\nL'entraide devient aventure."),
    ("🎯", "Pot Commun", "Financement solidaire des missions.\nContribue au bien commun."),
]

feat_w = Inches(3.8)
feat_h = Inches(1.4)
cols = 3
rows = 3
start_x = Inches(0.8)
start_y = Inches(1.8)

for i, (emoji, title, desc) in enumerate(features):
    col = i % cols
    row = i // cols
    x = start_x + col * (feat_w + Inches(0.3))
    y = start_y + row * (feat_h + Inches(0.2))
    
    add_shape(slide, x, y, feat_w, feat_h, fill_color=WHITE, line_color=BORDER_GRAY, line_width=Pt(0.5), corner_radius=0.06)
    
    # Emoji on left
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5),
                emoji, font_size=24)
    
    # Title + desc
    add_textbox(slide, x + Inches(0.7), y + Inches(0.15), feat_w - Inches(1.0), Inches(0.4),
                title, font_size=16, bold=True, color=TEXT_DARK)
    add_textbox(slide, x + Inches(0.7), y + Inches(0.55), feat_w - Inches(1.0), Inches(0.7),
                desc, font_size=12, color=TEXT_BODY)

# ══════════════════════════════════════════════
# SLIDE 5 — GAMIFICATION
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
            "DIFFÉRENTIATEUR CLÉ", font_size=11, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.8),
            "Gamification Hero :\nton entraide devient ton XP",
            font_size=40, bold=True, color=TEXT_DARK)

gamification = [
    ("⭐", "5 Niveaux", "Débutant → Héros Légendaire.\nChaque niveau débloque\ndes privilèges."),
    ("🎖️", "Badges", "Badges thématiques pour les\nactions spéciales et\nl'engagement durable."),
    ("📋", "Quêtes", "Défis quotidiens et missions\nspéciales pour maintenir\nl'engagement."),
]

for i, (emoji, title, desc) in enumerate(gamification):
    x = start_x + i * (card_w + card_gap)
    y = Inches(2.2)
    
    add_shape(slide, x, y, card_w, Inches(3.2), fill_color=LIGHT_GRAY, corner_radius=0.08)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.25), Inches(0.6), Inches(0.5),
                emoji, font_size=30)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.6), Inches(0.5),
                title, font_size=22, bold=True, color=TEAL_DARK)
    add_textbox(slide, x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.6), Inches(1.5),
                desc, font_size=14, color=TEXT_BODY)

# Key differentiator box
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.2), fill_color=TEAL_LIGHT, corner_radius=0.06)
add_textbox(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(1.0),
            "≠ Différence fondamentale : La plupart des time banks s'arrêtent au matching.\n"
            "TimeHeroes transforme chaque échange en expérience engageante — comme un RPG social\n"
            "où plus tu aides, plus tu montes en grade. Et ça devient un CV de l'engagement.",
            font_size=14, color=TEXT_BODY)

# ══════════════════════════════════════════════
# SLIDE 6 — USER JOURNEY
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, OFF_WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
            "PARCOURS UTILISATEUR", font_size=11, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7),
            "Du besoin au héros en 5 minutes",
            font_size=40, bold=True, color=TEXT_DARK)

steps_left = [
    ("1", "Je crée mon compte Hero", "Inscription rapide avec email.\nOu connexion démo 1-clic."),
    ("2", "Je propose ou je cherche", "Je publie ce que je sais faire,\nou je trouve quelqu'un pour m'aider."),
    ("3", "Je réserve en quelques clics", "Les TIME sont bloqués en escrow.\nLe héros est prévenu."),
]
steps_right = [
    ("4", "La mission se réalise", "Cours de guitare, aide déménagement,\nsoutien scolaire…"),
    ("5", "Je valide par QR/NFC", "Un scan. Un tag. Les TIME sont libérés.\nLes deux parties notent."),
    ("6", "Je gagne XP + badges", "Ma réputation grandit. Mon niveau\naugmente. Mon CV s'enrichit."),
]

def draw_steps(slide, steps, offset_x):
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.0) + i * Inches(1.7)
        
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
            offset_x + Inches(0.3), y + Inches(0.1), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = TEAL
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False
        
        # Text
        add_textbox(slide, offset_x + Inches(1.1), y, Inches(4.5), Inches(0.4),
                    title, font_size=18, bold=True, color=TEXT_DARK)
        add_textbox(slide, offset_x + Inches(1.1), y + Inches(0.45), Inches(4.5), Inches(0.6),
                    desc, font_size=14, color=TEXT_BODY)
        
        # Connector line (except last)
        if i < len(steps) - 1:
            add_shape(slide, offset_x + Inches(0.5), y + Inches(0.6), Pt(2), Inches(1.0), fill_color=TEAL)

draw_steps(slide, steps_left, Inches(0.5))
draw_steps(slide, steps_right, Inches(6.5))

# ══════════════════════════════════════════════
# SLIDE 7 — METRICS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
            "CHIFFRES & RÉALISATION", font_size=11, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7),
            "Du concept à la réalité en 2 mois",
            font_size=40, bold=True, color=TEXT_DARK)

# Big stats
stats = [
    ("12", "Lots fonctionnels\nlivrés"),
    ("22+", "Routes /\npages"),
    ("10", "Catégories de\nservices"),
    ("5", "Niveaux\nHero"),
]

for i, (num, label) in enumerate(stats):
    x = Inches(0.8) + i * Inches(3.1)
    y = Inches(1.9)
    add_shape(slide, x, y, Inches(2.8), Inches(2.0), fill_color=LIGHT_GRAY, corner_radius=0.08)
    add_textbox(slide, x, y + Inches(0.2), Inches(2.8), Inches(0.8),
                num, font_size=48, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + Inches(1.0), Inches(2.8), Inches(0.8),
                label, font_size=13, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# Detail boxes
details = [
    ("Lots livrés", "Auth & Wallet · Marketplace · Booking & Escrow\nRating & Reputation · Local Heroes · Urgent Help\nNFC Proof · Gamification · Demo Seed · Landing Page"),
    ("Screens principales", "Dashboard · Marketplace · Détail Mission · Booking\nWallet · Rewards · Profile · Urgent Help · Admin"),
    ("Stack", "Next.js 16 · TypeScript · Prisma 6 · SQLite\nTailwind v4 · NextAuth · Zod"),
]

for i, (title, content) in enumerate(details):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(4.3)
    add_shape(slide, x, y, Inches(3.8), Inches(2.5), fill_color=LIGHT_GRAY, corner_radius=0.08)
    add_shape(slide, x, y, Inches(3.8), Inches(0.5), fill_color=TEAL, corner_radius=0.08)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.4), Inches(0.4),
                title, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.6),
                content, font_size=13, color=TEXT_BODY)

# ══════════════════════════════════════════════
# SLIDE 8 — ARCHITECTURE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, OFF_WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
            "ARCHITECTURE & DÉPLOIEMENT", font_size=11, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7),
            "Production ready",
            font_size=40, bold=True, color=TEXT_DARK)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8), Inches(0.4),
            "Une stack moderne, un déploiement VPS réel",
            font_size=16, color=TEXT_MUTED)

# Tech stack cloud
stack_items = ["Next.js 16", "TypeScript", "Prisma 6", "SQLite", "Tailwind v4", "NextAuth", "Zod", "Lucide Icons", "bcryptjs"]
y = Inches(2.2)
start_x_stack = Inches(0.8)
for i, item in enumerate(stack_items):
    col = i % 3
    row = i // 3
    x = start_x_stack + col * Inches(1.35)
    yy = y + row * Inches(0.55)
    add_shape(slide, x, yy, Inches(1.2), Inches(0.4), fill_color=TEAL_LIGHT, line_color=TEAL, line_width=Pt(0.5), corner_radius=0.05)
    add_textbox(slide, x, yy + Inches(0.02), Inches(1.2), Inches(0.35),
                item, font_size=11, bold=False, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

# Architecture diagram
add_shape(slide, Inches(0.8), Inches(4.3), Inches(6.0), Inches(2.5), fill_color=WHITE, line_color=BORDER_GRAY, line_width=Pt(1), corner_radius=0.06)
lines = [
    "Client → HTTPS timeheroes.fr",
    "  → Caddy (SSL/TLS) → localhost:3096",
    "    → Next.js (programmatic server)",
    "      → Prisma ORM → SQLite",
]
y_pos = Inches(4.5)
for line in lines:
    indent = line.count("  ")
    add_textbox(slide, Inches(1.0) + Inches(0.3 * indent), y_pos, Inches(5.5), Inches(0.3),
                line.strip(), font_size=12, color=TEXT_BODY, font_name='Consolas')
    y_pos += Inches(0.45)

# Deploy info card
add_shape(slide, Inches(7.3), Inches(2.2), Inches(5.2), Inches(4.6), fill_color=WHITE, line_color=TEAL, line_width=Pt(1.5), corner_radius=0.08)
add_textbox(slide, Inches(7.6), Inches(2.5), Inches(4.6), Inches(0.5),
            "🌐 timeheroes.fr", font_size=22, bold=True, color=TEXT_DARK)
add_textbox(slide, Inches(7.6), Inches(3.1), Inches(4.6), Inches(0.5),
            "Hébergé sur VPS Hetzner (Allemagne)\nDomaine OVH · SSL Let's Encrypt",
            font_size=13, color=TEXT_BODY)

# Demo box
add_shape(slide, Inches(7.6), Inches(4.0), Inches(4.6), Inches(2.5), fill_color=TEAL_LIGHT, corner_radius=0.06)
add_textbox(slide, Inches(7.9), Inches(4.2), Inches(4.0), Inches(0.3),
            "🔑 TESTEZ LA DÉMO", font_size=11, bold=True, color=TEAL_DARK)
add_textbox(slide, Inches(7.9), Inches(4.6), Inches(4.0), Inches(0.8),
            "Email : demo@timeheroes.fr\nMot de passe : TimeHeroes2026!",
            font_size=15, bold=True, color=TEXT_DARK)
add_textbox(slide, Inches(7.9), Inches(5.5), Inches(4.0), Inches(0.5),
            "Ou utilisez ⚡ Connexion démo 1-clic sur la page de login",
            font_size=12, color=TEXT_BODY)

# ══════════════════════════════════════════════
# SLIDE 9 — ROADMAP
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, WHITE)
add_teal_bar(slide, height=Pt(4))

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
            "FEUILLE DE ROUTE", font_size=11, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.7),
            "La suite de l'aventure TimeHeroes",
            font_size=40, bold=True, color=TEXT_DARK)

phases = [
    ("✅ Fait — POC", "Base fonctionnelle",
     ["Auth, Wallet, Marketplace", "Booking & Escrow", "QR/NFC, Gamification", "Urgent Help, Local Heroes", "Landing page, seed data"],
     TEAL, "Juin 2026"),
    ("🎯 En cours — MVP", "Passage à l'échelle",
     ["Dashboard admin analytics", "Notifications temps réel", "Messagerie intégrée", "Onboarding enrichi", "PostgreSQL scale"],
     ACCENT_ORANGE, "Q3 2026"),
    ("🌟 Vision — Long terme", "Impact & écosystème",
     ["Records of Achievement", "Portfolio bénévole / CV", "Réseau coopératives ESS", "Application mobile", "API publique"],
     ACCENT_PURPLE, "2027"),
]

phase_w = Inches(3.6)
for i, (label, title, items, accent, period) in enumerate(phases):
    x = Inches(0.8) + i * (phase_w + Inches(0.4))
    y = Inches(2.0)
    
    # Card
    add_shape(slide, x, y, phase_w, Inches(4.8), fill_color=LIGHT_GRAY, corner_radius=0.08)
    
    # Accent top bar
    add_shape(slide, x, y, phase_w, Pt(5), fill_color=accent)
    
    # Period
    add_textbox(slide, x + Inches(0.3), y + Inches(0.2), phase_w - Inches(0.6), Inches(0.3),
                period, font_size=10, bold=True, color=TEXT_MUTED)
    
    # Label
    add_textbox(slide, x + Inches(0.3), y + Inches(0.5), phase_w - Inches(0.6), Inches(0.4),
                label, font_size=14, bold=True, color=accent)
    
    # Title
    add_textbox(slide, x + Inches(0.3), y + Inches(1.0), phase_w - Inches(0.6), Inches(0.5),
                title, font_size=20, bold=True, color=TEXT_DARK)
    
    # Items
    item_y = y + Inches(1.7)
    for item in items:
        add_textbox(slide, x + Inches(0.3), item_y, phase_w - Inches(0.6), Inches(0.35),
                    f"→  {item}", font_size=13, color=TEXT_BODY)
        item_y += Inches(0.45)

# ══════════════════════════════════════════════
# SLIDE 10 — VISION & CTA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_bg(slide, OFF_WHITE)
add_teal_bar(slide, height=Pt(8))

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(4), Inches(0.4),
            "🌍 L'AVENIR DE L'ENTRAIDE", font_size=12, bold=True, color=TEAL_DARK)

add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(1.5),
            "L'engagement citoyen\n",
            font_size=56, bold=True, color=TEXT_DARK)

add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.5), Inches(0.8),
            "devient un actif professionnel",
            font_size=56, bold=True, color=TEAL)

# Vision text
add_textbox(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(1.0),
            "TimeHeroes, c'est plus qu'une app. C'est un nouveau standard pour\n"
            "reconnaître, valoriser et monétiser l'engagement de chacun.\n"
            "Le temps est la seule vraie richesse. Échangeons-la.",
            font_size=18, color=TEXT_BODY)

# CTA
add_shape(slide, Inches(0.8), Inches(5.0), Inches(4.5), Inches(0.8), fill_color=TEAL, corner_radius=0.1)
add_textbox(slide, Inches(0.8), Inches(5.05), Inches(4.5), Inches(0.7),
            "🚀  timeheroes.fr", font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Demo credentials
add_shape(slide, Inches(0.8), Inches(6.0), Inches(5.5), Inches(0.6), fill_color=WHITE, line_color=BORDER_GRAY, line_width=Pt(1), corner_radius=0.06)
add_textbox(slide, Inches(1.0), Inches(6.05), Inches(5.0), Inches(0.5),
            "📧  demo@timeheroes.fr   /   TimeHeroes2026!",
            font_size=14, color=TEXT_BODY)

# Right decorative element
add_shape(slide, Inches(9.0), Inches(1.0), Inches(3.5), Inches(3.5), fill_color=TEAL_LIGHT, corner_radius=0.15)
add_textbox(slide, Inches(9.3), Inches(1.8), Inches(2.9), Inches(2.0),
            "PITCH\nDECK\n2026", font_size=40, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

# Presenter
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(5), Inches(0.4),
            "Ronald Mounien — MBA ESSEC Executive 2026",
            font_size=14, color=TEXT_MUTED)

# ── Save ──
output_path = "/root/projects/timebank-poc/TimeHeroes - MBA Pitch Deck.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
