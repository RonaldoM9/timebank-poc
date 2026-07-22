"""
Generate DEMO_P2P_FACILITATOR_TEST_PLAN.pptx
Test scenarios for TimeHeroes P2P & Facilitator features
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ─── Color Palette (Light theme TimeHeroes) ───────────────────────────────
BG        = RGBColor(0xF8, 0xF5, 0xF0)   # warm cream
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SURFACE2   = RGBColor(0xF0, 0xED, 0xE8)   # light gray card bg
ACCENT     = RGBColor(0x00, 0xD4, 0xAA)    # teal
ACCENT_DARK= RGBColor(0x00, 0xA8, 0x89)
TEXT       = RGBColor(0x10, 0x10, 0x10)
TEXT_SEC   = RGBColor(0x5F, 0x63, 0x68)
TEXT_MUTED = RGBColor(0x9A, 0x9E, 0xA2)
BORDER     = RGBColor(0xE0, 0xDD, 0xD8)
ORANGE     = RGBColor(0xF5, 0x92, 0x35)
RED        = RGBColor(0xE8, 0x4D, 0x50)
PURPLE     = RGBColor(0xAB, 0x47, 0xBC)
GREEN      = RGBColor(0x38, 0xC1, 0x72)

# ─── Helpers ──────────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color=None, border_color=None, border_width=None, radius=None):
    if radius:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        # Set corner radius
        sp = shape._element
        prstGeom = sp.find(qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                from lxml import etree
                avLst = etree.SubElement(prstGeom, qn('a:avLst'))
            from lxml import etree
            gd = etree.SubElement(avLst, qn('a:gd'))
            gd.set('name', 'adj')
            gd.set('fmla', f'val {int(radius)}')
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    if color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text="", font_size=14,
                 bold=False, color=TEXT, align=PP_ALIGN.LEFT, font_name='Calibri'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf

def add_rich_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tf

def add_para(tf, text, font_size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             space_after=6, font_name='Calibri', bullet=False):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(space_after)
    if bullet:
        p.level = 0
        # Add bullet
        pPr = p._pPr
        if pPr is None:
            from lxml import etree
            pPr = etree.SubElement(p._p, qn('a:pPr'))
        from lxml import etree
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '●')
    return p

def add_bullet(tf, text, font_size=13, color=TEXT_SEC, level=0, space_after=4):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    p.space_after = Pt(space_after)
    p.level = level
    return p

def add_run(p, text, bold=False, color=None, font_size=None, font_name='Calibri'):
    run = p.add_run()
    run.text = text
    if bold: run.font.bold = True
    if color: run.font.color.rgb = color
    if font_size: run.font.size = Pt(font_size)
    run.font.name = font_name
    return run

def add_phase_bar(slide, left, top, width, color):
    """Colored bar accent at top of card"""
    bar = add_shape(slide, left, top, width, Pt(4), color)

def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Accent bar top
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT)
    # Title
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                 title, 36, True, TEXT, PP_ALIGN.LEFT, 'Calibri')
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3), Inches(11), Inches(1),
                     subtitle, 18, False, TEXT_SEC, PP_ALIGN.LEFT)
    # Teal circle decoration
    add_shape(slide, Inches(11.5), Inches(0.5), Inches(1.2), Inches(1.2),
              ACCENT, radius=60000)
    return slide

def add_section_slide(prs, number, title, description=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    # Phase bar top
    colors = [ACCENT, ORANGE, PURPLE, ACCENT_DARK, GREEN, ACCENT]
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), colors[number % len(colors)])
    # Number circle
    circ = add_shape(slide, Inches(1), Inches(1.2), Inches(0.9), Inches(0.9), ACCENT, radius=45000)
    tf = circ.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{number}"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_text_box(slide, Inches(2.2), Inches(1.2), Inches(10), Inches(0.8),
                 title, 28, True, TEXT, PP_ALIGN.LEFT)
    if description:
        add_text_box(slide, Inches(2.2), Inches(2.1), Inches(9), Inches(0.6),
                     description, 14, False, TEXT_SEC, PP_ALIGN.LEFT)
    return slide

def add_test_card(slide, left, top, width, height, test_id, title, steps, tag=None, tag_color=None):
    """A test scenario card"""
    card = add_shape(slide, left, top, width, height, WHITE, BORDER, 0.5, radius=30000)
    # Left accent
    add_shape(slide, left + Pt(1), top + Inches(0.15), Pt(3), height - Inches(0.3), tag_color or ACCENT)
    # Test ID badge
    badge = add_shape(slide, left + Inches(0.2), top + Inches(0.15), Inches(1), Inches(0.3),
                      tag_color or ACCENT, radius=15000)
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = test_id
    p.font.size = Pt(9)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), Inches(0.35),
                 title, 12, True, TEXT, PP_ALIGN.LEFT)
    # Steps
    y = top + Inches(0.95)
    tf = add_rich_textbox(slide, left + Inches(0.2), y, width - Inches(0.4), height - Inches(1.15))
    for i, step in enumerate(steps[:4]):
        add_bullet(tf, f"{step}", 9, TEXT_SEC, 0, 2)
    if len(steps) > 4:
        add_bullet(tf, f"+{len(steps)-4} étapes...", 9, TEXT_MUTED, 0, 2)
    return card

# ─── Build Presentation ────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ════════════════ SLIDE 1 — TITLE ════════════════
slide = add_title_slide(prs, "🧪 Plan de Test P2P & Facilitateur",
                        "TimeHeroes — Parcours pair-à-pair et outillage facilitateur")
# Credentials box
box = add_shape(slide, Inches(1), Inches(4.2), Inches(11), Inches(2.2), WHITE, ACCENT, 1.5, radius=30000)
tf = add_rich_textbox(slide, Inches(1.3), Inches(4.4), Inches(10.5), Inches(1.8))
add_para(tf, "🔑 Comptes de test", 16, True, ACCENT, space_after=8)
add_bullet(tf, "Alex Demo (ADMIN)  →  demo@timeheroes.fr  /  TimeHeroes2026!", 12, TEXT_SEC)
add_bullet(tf, "Sarah Martin (FACILITATOR)  →  sarah.demo@timeheroes.fr  /  TimeHeroes2026!", 12, TEXT_SEC)
add_bullet(tf, "Karim Benali (USER)  →  karim.demo@timeheroes.fr  /  TimeHeroes2026!", 12, TEXT_SEC)
add_bullet(tf, "Alice Seed (USER)  →  alice.seed@timeheroes.fr  /  TimeHeroes2026!", 12, TEXT_SEC)

# ════════════════ SLIDE 2 — SOMMAIRE ════════════════
slide = add_section_slide(prs, 1, "Sommaire", "30 scénarios répartis en 6 parties")
parts = [
    ("Partie 1 — Parcours P2P Complet", "S1 à S9", ACCENT),
    ("Partie 2 — Urgences", "S10 à S11", ORANGE),
    ("Partie 3 — Pot Commun", "S12 à S16", PURPLE),
    ("Partie 4 — Matching", "S17 à S21", ACCENT_DARK),
    ("Partie 5 — Intelligence Réseau", "S22 à S28", GREEN),
    ("Partie 6 — Missions Collectives", "S29 à S30", ACCENT),
]
for i, (name, ids, color) in enumerate(parts):
    y = Inches(3) + Inches(0.65) * i
    # Phase dot
    add_shape(slide, Inches(1.5), y + Inches(0.08), Inches(0.15), Inches(0.15), color, radius=7500)
    add_text_box(slide, Inches(1.9), y, Inches(6), Inches(0.3), name, 14, True, TEXT, PP_ALIGN.LEFT)
    add_text_box(slide, Inches(8.5), y, Inches(2), Inches(0.3), ids, 12, False, TEXT_MUTED, PP_ALIGN.RIGHT)

# ════════════════ SLIDE 3 — PRÉPARATION ════════════════
slide = add_section_slide(prs, 2, "Préparation", "Données de démo et réinitialisation")
tf = add_rich_textbox(slide, Inches(1.3), Inches(3.3), Inches(5.5), Inches(3.5))
add_para(tf, "🔄 Réinitialisation des données", 16, True, ACCENT, space_after=8)
add_bullet(tf, "npm run db:reset-demo", 12, TEXT_SEC)
add_bullet(tf, "ou lancer les seeds un par un :", 12, TEXT_SEC)
add_bullet(tf, "npm run seed:demo", 11, TEXT_MUTED, 1)
add_bullet(tf, "npm run seed:gamification", 11, TEXT_MUTED, 1)
add_bullet(tf, "npm run seed:facilitator-test", 11, TEXT_MUTED, 1)
add_bullet(tf, "npm run seed:collective-missions", 11, TEXT_MUTED, 1)
add_bullet(tf, "npm run seed:facilitator-network", 11, TEXT_MUTED, 1)

tf2 = add_rich_textbox(slide, Inches(7.3), Inches(3.3), Inches(5), Inches(3.5))
add_para(tf2, "📊 État actuel de la base", 16, True, ACCENT, space_after=8)
stats = [("Utilisateurs", "23"), ("Services", "47"), ("Réservations", "43"),
         ("Transactions", "75"), ("Évaluations", "20"), ("Demandes Pot", "6"),
         ("Missions collectives", "5"), ("Demandes urgentes", "9"),
         ("Messages P2P", "13"), ("Recommandations match", "54")]
for label, val in stats:
    row = tf2.add_paragraph()
    row.text = f"{label}  "
    row.font.size = Pt(11)
    row.font.color.rgb = TEXT_SEC
    run = row.add_run()
    run.text = val
    run.font.bold = True
    run.font.size = Pt(11)
    run.font.color.rgb = ACCENT
    row.space_after = Pt(3)

# ════════════════ PARTIE 1 — P2P ════════════════
slide = add_section_slide(prs, 3, "Partie 1 — Parcours P2P Complet",
                          "S1 à S9 : Marketplace, réservation, messagerie, complétion, wallet, transfert")

# S1-S3 row
add_test_card(slide, Inches(0.5), Inches(3.2), Inches(3.8), Inches(2.4),
              "S1", "Marketplace", [
                  "Connecté en Karim → /services",
                  "✅ Liste des services visible",
                  "✅ Cartes : nom, note, TIME/h",
                  "✅ Page détail service (description, profil, avis)"
              ], "P2P", ACCENT)

add_test_card(slide, Inches(4.6), Inches(3.2), Inches(3.8), Inches(2.4),
              "S2", "Réservation P2P", [
                  "Clique Réserver sur un service",
                  "✅ Remplis heures + date",
                  "✅ Récap coût (heures × TIME/h)",
                  "✅ Statut : En attente (pending)"
              ], "P2P", ACCENT)

add_test_card(slide, Inches(8.7), Inches(3.2), Inches(4.1), Inches(2.4),
              "S3", "Messagerie P2P", [
                  "Va sur Mes réservations → /bookings",
                  "✅ Envoie message au provider",
                  "✅ Le message s'affiche avec nom/date",
                  "✅ Discussion visible dans le détail"
              ], "P2P", ACCENT)

# S4-S6 row
add_test_card(slide, Inches(0.5), Inches(5.8), Inches(3.8), Inches(1.5),
              "S4", "Réponse du Provider", [
                  "Déconnecte-toi, connecte le provider",
                  "✅ Répond dans la discussion",
                  "✅ Marqueur de lecture (read)"
              ], "P2P", ORANGE)

add_test_card(slide, Inches(4.6), Inches(5.8), Inches(3.8), Inches(1.5),
              "S5", "Changement de statut", [
                  "✅ Provider clique Accepter",
                  "✅ Statut → Confirmé",
                  "✅ Message système dans la discussion"
              ], "P2P", ORANGE)

add_test_card(slide, Inches(8.7), Inches(5.8), Inches(4.1), Inches(1.5),
              "S6", "Complétion (QR)", [
                  "✅ Générer QR de complétion",
                  "✅ Scanner / copier lien autre session",
                  "✅ Statut → Terminée"
              ], "P2P", ORANGE)

# S7-S9 next slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(6), Inches(0.5),
             "Partie 1 — P2P (suite)", 20, True, TEXT, PP_ALIGN.LEFT)

add_test_card(slide, Inches(0.5), Inches(1.5), Inches(3.8), Inches(2.8),
              "S7", "Évaluation", [
                  "✅ Noter l'expérience (1-5★)",
                  "✅ Commentaire",
                  "✅ Note visible sur profil provider",
                  "✅ Réputation impactée"
              ], "P2P", ACCENT)

add_test_card(slide, Inches(4.6), Inches(1.5), Inches(3.8), Inches(2.8),
              "S8", "Wallet & Transactions", [
                  "✅ Solde débité du montant",
                  "✅ escrow_hold dans historique",
                  "✅ escrow_release au provider",
                  "✅ Solde provider augmenté"
              ], "P2P", PURPLE)

add_test_card(slide, Inches(8.7), Inches(1.5), Inches(4.1), Inches(2.8),
              "S9", "Transfert P2P direct", [
                  "Va sur /wallet/transfer",
                  "✅ Envoie 1 TIME à alice.seed@...",
                  "✅ Message de confirmation",
                  "✅ Transaction dans les 2 historiques"
              ], "P2P", PURPLE)

# ════════════════ PARTIE 2 — URGENCES ════════════════
slide = add_section_slide(prs, 4, "Partie 2 — Urgences",
                          "S10 à S11 : Création et réponse aux demandes urgentes")

add_test_card(slide, Inches(2.5), Inches(3.2), Inches(3.8), Inches(3.5),
              "S10", "Créer une demande urgente", [
                  "Connecté en Karim → /urgent",
                  "✅ Liste des demandes existantes",
                  "✅ Nouvelle demande : titre, desc, catégorie",
                  "✅ Soumettre",
                  "✅ Statut : Ouvert"
              ], "URGENCE", ORANGE)

add_test_card(slide, Inches(6.8), Inches(3.2), Inches(3.8), Inches(3.5),
              "S11", "Répondre à une urgence", [
                  "Connecte-toi en Alex → /urgent",
                  "✅ Demande de Karim visible",
                  "✅ Proposer son aide → envoyer offre",
                  "✅ Reviens en Karim → offre visible",
                  "✅ Accepter l'offre → statut Résolu"
              ], "URGENCE", ORANGE)

# ════════════════ PARTIE 3 — POT COMMUN ════════════════
slide = add_section_slide(prs, 5, "Partie 3 — Pot Commun",
                          "S12 à S16 : Dashboard, approbation/refus, transactions, missions solidaires")

add_test_card(slide, Inches(0.5), Inches(3.2), Inches(3.8), Inches(1.8),
              "S12", "Dashboard Pot Commun", [
                  "Connecté en Sarah → /facilitator/community-pot",
                  "✅ Solde pot, donations, financements",
                  "✅ Demandes en attente, KPIs interactifs"
              ], "POT", PURPLE)

add_test_card(slide, Inches(4.6), Inches(3.2), Inches(3.8), Inches(1.8),
              "S13", "Approuver une demande", [
                  "✅ Vérifier détails demande",
                  "✅ Approuver avec note de décision",
                  "✅ Solde pot diminue"
              ], "POT", PURPLE)

add_test_card(slide, Inches(8.7), Inches(3.2), Inches(4.1), Inches(1.8),
              "S14", "Refuser une demande", [
                  "✅ Refuser avec note",
                  "✅ Passe en Refusée",
                  "✅ Solde pot inchangé"
              ], "POT", PURPLE)

add_test_card(slide, Inches(0.5), Inches(5.3), Inches(3.8), Inches(1.8),
              "S15", "Transactions du pot", [
                  "✅ Dernières transactions visibles",
                  "✅ Donations (entrées) + Fundings (sorties)",
                  "✅ Montant, user, date, raison"
              ], "POT", ACCENT_DARK)

add_test_card(slide, Inches(4.6), Inches(5.3), Inches(3.8), Inches(1.8),
              "S16", "Missions solidaires", [
                  "✅ Vérifier missions SELF_DECLARED",
                  "✅ Action Vérifier",
                  "✅ Action Refuser"
              ], "SOLIDAIRE", ACCENT)

# ════════════════ PARTIE 4 — MATCHING ════════════════
slide = add_section_slide(prs, 6, "Partie 4 — Matching",
                          "S17 à S21 : Recommandations pour urgences, solidaires, collectives, historique")

add_test_card(slide, Inches(0.5), Inches(3.2), Inches(3.8), Inches(2.0),
              "S17", "Matching urgent", [
                  "Va sur /facilitator/matching",
                  "✅ Onglet Urgent",
                  "✅ Générer recommandations",
                  "✅ Heroes recommandés avec score"
              ], "MATCH", GREEN)

add_test_card(slide, Inches(4.6), Inches(3.2), Inches(3.8), Inches(2.0),
              "S18", "Feedback matching", [
                  "✅ Approuver / Rejeter recommandation",
                  "✅ Laisser feedback",
                  "✅ Feedback enregistré"
              ], "MATCH", GREEN)

add_test_card(slide, Inches(8.7), Inches(3.2), Inches(4.1), Inches(2.0),
              "S19", "Matching solidaire", [
                  "✅ Onglet Solidaire",
                  "✅ Générer recommandations",
                  "✅ Résultats affichés"
              ], "MATCH", ACCENT_DARK)

add_test_card(slide, Inches(0.5), Inches(5.5), Inches(3.8), Inches(1.8),
              "S20", "Matching collectif", [
                  "✅ Onglet Collective",
                  "✅ Générer recommandations",
                  "✅ Résultats affichés"
              ], "MATCH", ACCENT_DARK)

add_test_card(slide, Inches(4.6), Inches(5.5), Inches(3.8), Inches(1.8),
              "S21", "Historique matchs", [
                  "✅ Onglet Historique",
                  "✅ Recommandations passées visibles",
                  "✅ Feedbacks associés"
              ], "MATCH", ACCENT)

# ════════════════ PARTIE 5 — INTELLIGENCE RÉSEAU ════════════════
slide = add_section_slide(prs, 7, "Partie 5 — Intelligence Réseau",
                          "S22 à S28 : Dashboard réseau, demandes bloquées, héros, TIME dormants, alertes")

add_test_card(slide, Inches(0.5), Inches(3.2), Inches(3.8), Inches(2.2),
              "S22", "Dashboard Réseau", [
                  "Va sur /facilitator/network",
                  "✅ Network Health Score (0-100)",
                  "✅ 5 sous-scores : liquidité, réponse, etc.",
                  "✅ Métriques clés affichées"
              ], "RÉSEAU", GREEN)

add_test_card(slide, Inches(4.6), Inches(3.2), Inches(3.8), Inches(2.2),
              "S23", "Demandes bloquées", [
                  "✅ Urgences sans réponse",
                  "✅ Pot en attente, solidaires non vérifiées",
                  "✅ Bookings bloqués, collectives sous-rempl.",
                  "✅ Code couleur par sévérité"
              ], "RÉSEAU", ORANGE)

add_test_card(slide, Inches(8.7), Inches(3.2), Inches(4.1), Inches(2.2),
              "S24", "Heroes sur-sollicités", [
                  "✅ Missions 30j, heures données",
                  "✅ Score + niveau de risque",
                  "✅ Recommandations d'action",
                  "✅ Lien vers profil"
              ], "RÉSEAU", RED)

add_test_card(slide, Inches(0.5), Inches(5.7), Inches(3.8), Inches(1.5),
              "S25", "Heroes sous-utilisés", [
                  "✅ Compétences, passport, inactivité",
                  "✅ Recommandations d'action"
              ], "RÉSEAU", ACCENT)

add_test_card(slide, Inches(4.6), Inches(5.7), Inches(3.8), Inches(1.5),
              "S26", "TIME dormants", [
                  "✅ Statuts : Light, Strong, TimeRich, Poor",
                  "✅ Suggestions recyclage"
              ], "RÉSEAU", PURPLE)

add_test_card(slide, Inches(8.7), Inches(5.7), Inches(4.1), Inches(1.5),
              "S27", "Alertes & Notes", [
                  "✅ Résoudre / Ignorer une alerte",
                  "✅ Ajouter note sur hero ou mission",
                  "✅ Note visible dans historique"
              ], "RÉSEAU", ORANGE)

# S28
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Pt(4), ACCENT)
add_text_box(slide, Inches(0.5), Inches(0.3), Inches(6), Inches(0.5),
             "Partie 5 — Intelligence Réseau (suite)", 20, True, TEXT, PP_ALIGN.LEFT)

add_test_card(slide, Inches(1.5), Inches(1.5), Inches(3.8), Inches(2.0),
              "S28", "Rafraîchissement", [
                  "✅ Clique sur Rafraîchir",
                  "✅ Scores recalculés",
                  "✅ Données mises à jour"
              ], "RÉSEAU", GREEN)

# ════════════════ PARTIE 6 — MISSIONS COLLECTIVES ════════════════
slide = add_section_slide(prs, 8, "Partie 6 — Missions Collectives",
                          "S29 à S30 : Découverte et participation")

add_test_card(slide, Inches(1.5), Inches(3.2), Inches(3.8), Inches(2.5),
              "S29", "Voir les missions collectives", [
                  "Va sur /collective-missions",
                  "✅ Liste des missions",
                  "✅ Infos : titre, organisateur, participants"
              ], "COLLECTIF", ACCENT)

add_test_card(slide, Inches(6.8), Inches(3.2), Inches(3.8), Inches(2.5),
              "S30", "Participer à une mission", [
                  "✅ Page détail : description, créneau",
                  "✅ Type : ONE_TO_MANY, MANY_TO_ONE...",
                  "✅ Clique Participer → compteur +1"
              ], "COLLECTIF", PURPLE)

# ════════════════ SLIDE FINALE — SCÉNARIO COMPLET ════════════════
slide = add_title_slide(prs, "🎯 Scénario Complet (Journey)",
                        "Enchaînement réaliste pour un facilitateur")
tf = add_rich_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(3.5))
steps = [
    ("Matin", "Sarah se connecte, vérifie le Dashboard Réseau", ACCENT),
    ("🔍", "Repère une urgence bloquée depuis 48h", ORANGE),
    ("🎯", "Matching → génère recommandations → approuve un Hero", GREEN),
    ("🏺", "Pot Commun → approuve une demande de financement", PURPLE),
    ("📝", "Laisse une note sur un Hero pour suivi", ACCENT_DARK),
    ("⏰", "Vérifie les TIME dormants et suggestions de recyclage", ORANGE),
    ("🔄", "Clique Rafraîchir → nouveau score réseau", ACCENT),
]
for i, (icon, text, color) in enumerate(steps):
    y = Inches(3.6) + Inches(0.6) * i
    add_shape(slide, Inches(1.2), y + Inches(0.05), Inches(0.15), Inches(0.15), color, radius=7500)
    add_text_box(slide, Inches(1.5), y, Inches(0.4), Inches(0.3), icon, 12, False, TEXT, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.2), y, Inches(9), Inches(0.3), text, 14, False, TEXT_SEC, PP_ALIGN.LEFT)

# ════════════════ CHECKLIST ════════════════
slide = add_title_slide(prs, "✅ Checklist finale",
                        "30 tests à cocher au fur et à mesure")
tf = add_rich_textbox(slide, Inches(1), Inches(3), Inches(11), Inches(4))
sections = [
    ("Partie 1 — P2P Complet", ["S1", "S2", "S3", "S4", "S5", "S6", "S7", "S8", "S9"]),
    ("Partie 2 — Urgences", ["S10", "S11"]),
    ("Partie 3 — Pot Commun", ["S12", "S13", "S14", "S15", "S16"]),
    ("Partie 4 — Matching", ["S17", "S18", "S19", "S20", "S21"]),
    ("Partie 5 — Intelligence Réseau", ["S22", "S23", "S24", "S25", "S26", "S27", "S28"]),
    ("Partie 6 — Missions Collectives", ["S29", "S30"]),
]
for i, (name, ids) in enumerate(sections):
    add_para(tf, name, 14, True, ACCENT, space_after=4)
    line = "  ".join([f"⬜ {tid}" for tid in ids])
    add_para(tf, line, 11, False, TEXT_SEC, space_after=10)

# ════════════════ SAVE ════════════════
output_path = "/root/projects/timebank-poc/DEMO_P2P_FACILITATOR_TEST_PLAN.pptx"
prs.save(output_path)
print(f"✅ Saved to {output_path}")
print(f"📊 Slides: {len(prs.slides)}")
